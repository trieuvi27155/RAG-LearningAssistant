"""Sinh bộ tài liệu mẫu ĐA DẠNG vào data/raw/ để đo chất lượng hệ thống.

Vì sao cần script này thay vì chỉ dùng 1 tài liệu thật:
Toàn bộ hiệu chỉnh trước đây của hệ thống (ngưỡng, chunk size, cách gộp ngữ cảnh) đều đo
trên ĐÚNG MỘT tài liệu - giáo trình PDF 230 trang thuần văn bản. Điều đó khiến mọi con số
"đã đo" chỉ chắc chắn đúng cho đúng loại tài liệu đó. Muốn biết hệ thống có thật sự tổng
quát hay không thì phải có tài liệu NGẮN, tài liệu nhiều BẢNG, tài liệu có HÌNH, tài liệu
nhiều cấp TIÊU ĐỀ - và phải có ngay để đo được trước/sau mỗi thay đổi.

Bộ tài liệu này cố tình cài sẵn những "bẫy" đã biết là chỗ code dễ sai:
  - Bảng lồng trong ô của bảng khác (python-docx không nhìn thấy - giới hạn đã ghi nhận).
  - Ảnh/bảng nằm trong GROUP SHAPE của PowerPoint (vòng lặp shape phẳng sẽ bỏ sót).
  - Ảnh neo/nổi trong DOCX (document.inline_shapes bỏ sót, phải duyệt qua rels).
  - Tiêu đề PDF chỉ phân biệt được bằng CỠ CHỮ (không có metadata cấu trúc như DOCX/PPTX).
  - Trang thưa chữ (slide) vs trang dày chữ - 2 cực của bài toán dựng ngữ cảnh.

LƯU Ý QUAN TRỌNG: tài liệu tổng hợp KHÔNG thay thế được tài liệu thật. Số đo trên đây là
chỉ dấu để phát hiện hồi quy và so sánh tương đối giữa các phiên bản, không phải bằng chứng
hệ thống chạy tốt trên tài liệu thật của người dùng. Mọi kết luận rút ra từ riêng bộ này
phải được đo lại khi có tài liệu thật.

Cách chạy:  python evaluation/tao_tai_lieu_mau.py
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from docx import Document
from docx.enum.text import WD_BREAK
from docx.shared import Inches as DocxInches
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

import config

THU_MUC_ANH_TAM = config.DATA_DIR / "anh_mau"

# Font mặc định của reportlab (Helvetica) KHÔNG có glyph tiếng Việt có dấu - chữ sẽ ra ô
# vuông hoặc mất dấu, khiến tài liệu mẫu vô dụng cho việc đo tiếng Việt. Dùng font hệ thống.
_FONT_HE_THONG = Path("C:/Windows/Fonts/arial.ttf")
_TEN_FONT = "ArialVN"


# ============================================================
# Nội dung mẫu (chủ đề trung tính, KHÔNG lấy từ tài liệu thật của người dùng)
# ============================================================
# Cố ý chọn chủ đề "quản lý thư viện" - đủ xa lĩnh vực pháp luật của tài liệu cũ để nếu
# model tình cờ nhớ nội dung cũ thì cũng không giúp gì được, và đủ đời thường để người đọc
# tự kiểm chứng câu trả lời đúng/sai mà không cần chuyên môn.

# Mỗi chương = 1 trang (ngăn bằng ngắt trang cứng). Nội dung từng chương phải KHÁC HẲN
# nhau: nếu các trang giống nhau thì Precision@K mất ý nghĩa (không phân biệt nổi hệ thống
# lấy đúng trang hay lấy trang bất kỳ có nội dung trùng) - đây là lỗi đã gặp ở bản đầu của
# chính script này, phát hiện khi soi lại nội dung từng trang.
CAC_CHUONG_DAI = [
    [
        ("Chương 1. Tổng quan về hệ thống thư viện", 1),
        ("1.1. Khái niệm thư viện số", 2),
        ("Thư viện số là hệ thống lưu trữ và cung cấp tài liệu dưới dạng điện tử, cho phép "
         "người đọc tra cứu và mượn tài liệu từ xa mà không cần đến trực tiếp. Khác với thư "
         "viện truyền thống vốn bị giới hạn bởi không gian vật lý và số bản in, thư viện số "
         "có thể phục vụ nhiều người đọc cùng lúc trên cùng một đầu tài liệu. Chi phí vận "
         "hành chủ yếu nằm ở hạ tầng máy chủ và bản quyền nội dung, không phải ở diện tích "
         "kho và nhân sự trông coi.", 0),
        ("1.2. Lịch sử hình thành", 2),
        ("Thư viện số đầu tiên được xây dựng vào năm 1971 với mục tiêu số hoá các tác phẩm "
         "đã hết hạn bản quyền. Đến thập niên 1990, khi Internet phổ biến, mô hình này mới "
         "thực sự lan rộng. Giai đoạn 2000 đến 2010 chứng kiến sự bùng nổ của các kho lưu "
         "trữ học thuật mở, cho phép công bố nghiên cứu miễn phí tới người đọc.", 0),
        ("1.3. Phân loại thư viện", 2),
        ("Thư viện được phân thành bốn nhóm chính theo đối tượng phục vụ: thư viện công "
         "cộng phục vụ mọi tầng lớp dân cư, thư viện chuyên ngành phục vụ một lĩnh vực hẹp, "
         "thư viện trường học phục vụ giảng dạy, và thư viện quốc gia giữ vai trò lưu chiểu "
         "toàn bộ xuất bản phẩm trong nước.", 0),
    ],
    [
        ("Chương 2. Nghiệp vụ quản lý mượn trả", 1),
        ("2.1. Quy trình mượn tài liệu", 2),
        ("Người đọc muốn mượn tài liệu phải có thẻ thư viện còn hiệu lực. Thủ thư kiểm tra "
         "tình trạng thẻ, số tài liệu đang mượn và các khoản phạt chưa thanh toán trước khi "
         "cho mượn thêm. Mỗi thẻ thường được mượn tối đa năm tài liệu cùng lúc trong thời "
         "hạn hai mươi mốt ngày.", 0),
        ("2.2. Quy trình trả và gia hạn", 2),
        ("Tài liệu trả quá hạn bị tính phí theo ngày. Người đọc có thể gia hạn trực tuyến "
         "tối đa hai lần nếu tài liệu đó không có người khác đặt trước. Trường hợp làm mất "
         "hoặc hư hỏng nặng, người đọc phải đền bù theo giá bìa cộng phí xử lý nghiệp vụ.", 0),
        ("2.3. Đặt trước tài liệu", 2),
        ("Khi tài liệu đang có người mượn, người đọc có thể đăng ký đặt trước để được ưu "
         "tiên nhận ngay sau khi tài liệu được trả về. Hàng đợi đặt trước xử lý theo thứ tự "
         "đăng ký. Người đặt trước có ba ngày để đến nhận kể từ khi nhận thông báo, quá hạn "
         "thì lượt ưu tiên chuyển cho người kế tiếp.", 0),
    ],
    [
        ("Chương 3. Bảo quản và số hoá", 1),
        ("3.1. Điều kiện bảo quản kho sách", 2),
        ("Kho sách cần duy trì nhiệt độ ổn định và độ ẩm trong ngưỡng cho phép để hạn chế "
         "nấm mốc và côn trùng. Ánh sáng trực tiếp làm giấy ố vàng và mực phai nhanh, nên "
         "khu vực lưu trữ lâu dài thường không có cửa sổ. Việc kiểm kê định kỳ giúp phát "
         "hiện sớm tài liệu xuống cấp cần đưa đi phục chế.", 0),
        ("3.2. Quy trình số hoá tài liệu cũ", 2),
        ("Tài liệu quý hiếm được quét ở độ phân giải cao rồi lưu ở định dạng không nén để "
         "giữ nguyên chi tiết. Bản quét sau đó được nhận dạng ký tự quang học nhằm tạo lớp "
         "văn bản tìm kiếm được. Bản gốc được đưa trở lại kho bảo quản và hạn chế lấy ra.", 0),
        ("3.3. Phục chế tài liệu hư hỏng", 2),
        ("Tài liệu rách hoặc bong gáy được chuyển tới bộ phận phục chế để gia cố bằng giấy "
         "dó và hồ trung tính. Không dùng băng dính thông thường vì lớp keo của nó ngả vàng "
         "và ăn mòn giấy sau vài năm. Mỗi lần phục chế đều được ghi nhật ký để theo dõi "
         "tình trạng tài liệu qua thời gian.", 0),
    ],
    [
        ("Chương 4. Thống kê và báo cáo", 1),
        ("4.1. Các chỉ số theo dõi", 2),
        ("Thư viện theo dõi bốn chỉ số chính: tổng lượt mượn, số thẻ hoạt động, tỉ lệ tài "
         "liệu trả đúng hạn, và thời gian chờ trung bình của hàng đợi đặt trước. Các chỉ số "
         "này được tổng hợp hàng tháng và so sánh với cùng kỳ năm trước.", 0),
        ("4.2. Kế hoạch bổ sung tài liệu", 2),
        ("Đầu sách có tỉ lệ đặt trước cao liên tục trong ba tháng được đề xuất mua thêm bản. "
         "Ngược lại, tài liệu không phát sinh lượt mượn nào trong hai năm sẽ được đưa vào "
         "danh sách rà soát thanh lý nhằm giải phóng diện tích kho.", 0),
    ],
]

# Tài liệu "gây nhiễu": rất nhiều mục CÙNG chủ đề, CÙNG cách diễn đạt, chỉ khác chi tiết
# (mã loại, số ngày, tên phòng, mức phí). Đây mới là dạng khó thật sự của tài liệu dài -
# và là dạng mà bộ mẫu ban đầu KHÔNG có: khi mỗi tài liệu nói một chủ đề riêng với từ vựng
# riêng thì tìm kiếm vector trúng đích ngay từ hạng 1, không phân biệt được hệ thống tốt hay
# tệ (đã đo: MRR = 1.00 ở mọi câu, tức phép đo bão hoà, vô dụng để so sánh).
#
# Với các mục gần như trùng nhau, việc tìm đúng mục phụ thuộc vào một vài TỪ KHOÁ HIẾM (mã
# hồ sơ, con số) - đúng chỗ tìm kiếm vector yếu và BM25/rerank mạnh.
_MA_HO_SO = ["A1", "A2", "B1", "B2", "C1", "C2", "D1", "D2", "E1", "E2",
             "F1", "F2", "G1", "G2", "H1", "H2", "K1", "K2", "M1", "M2"]
_PHONG = ["Phòng Nghiệp vụ", "Phòng Bạn đọc", "Phòng Địa chí", "Phòng Ngoại văn",
          "Phòng Thiếu nhi", "Phòng Tin học", "Phòng Bổ sung", "Phòng Bảo quản",
          "Phòng Hành chính", "Phòng Số hoá"]


def _sinh_muc_tuong_tu():
    """Sinh các mục quy định gần giống nhau, chỉ khác chi tiết - nguồn gây nhiễu có kiểm soát."""
    cac_muc = []
    for i, ma in enumerate(_MA_HO_SO):
        so_ngay = 3 + (i * 2) % 25
        phi = 10 + (i * 7) % 90
        phong = _PHONG[i % len(_PHONG)]
        so_ban = 1 + (i * 3) % 9
        cac_muc.append((f"Điều {i + 1}. Quy định xử lý hồ sơ loại {ma}", 2))
        cac_muc.append(
            (f"Hồ sơ loại {ma} được tiếp nhận tại {phong} và phải xử lý xong trong "
             f"{so_ngay} ngày làm việc kể từ ngày nhận đủ giấy tờ hợp lệ. Mức phí xử lý hồ "
             f"sơ loại {ma} là {phi} nghìn đồng. Mỗi lượt nộp được đăng ký tối đa {so_ban} "
             f"bản sao. Trường hợp hồ sơ loại {ma} thiếu giấy tờ, {phong} gửi thông báo bổ "
             f"sung trong vòng hai ngày làm việc và thời hạn xử lý được tính lại từ đầu.", 0)
        )
    return cac_muc


NOI_DUNG_SLIDE = [
    ("Hệ thống quản lý thư viện", ["Giới thiệu tổng quan", "Trình bày nội bộ"]),
    ("Mục tiêu hệ thống", ["Quản lý mượn trả tự động", "Tra cứu trực tuyến",
                           "Thống kê sử dụng"]),
    ("Đối tượng sử dụng", ["Người đọc", "Thủ thư", "Quản trị hệ thống"]),
    ("Quy tắc mượn", ["Tối đa 5 tài liệu mỗi thẻ", "Thời hạn 21 ngày",
                      "Gia hạn tối đa 2 lần"]),
    ("Phí phạt quá hạn", ["Tính theo ngày", "Miễn phạt 3 ngày đầu"]),
    ("Bảo quản kho", ["Kiểm soát nhiệt độ và độ ẩm", "Tránh ánh sáng trực tiếp"]),
    ("Số hoá tài liệu", ["Quét độ phân giải cao", "Nhận dạng ký tự quang học"]),
    ("Kết luận", ["Hệ thống giảm tải công việc thủ công"]),
]

BANG_PHI_PHAT = [
    ["Loại tài liệu", "Thời hạn mượn", "Phí phạt mỗi ngày", "Số lần gia hạn"],
    ["Sách giáo trình", "21 ngày", "2.000 đồng", "2 lần"],
    ["Sách tham khảo", "14 ngày", "3.000 đồng", "1 lần"],
    ["Tạp chí đóng tập", "7 ngày", "5.000 đồng", "Không gia hạn"],
    ["Tài liệu quý hiếm", "Đọc tại chỗ", "Không áp dụng", "Không gia hạn"],
]

BANG_PHAN_QUYEN = [
    ["Vai trò", "Xem tài liệu", "Mượn tài liệu", "Sửa dữ liệu"],
    ["Người đọc", "Có", "Có", "Không"],
    ["Thủ thư", "Có", "Có", "Có"],
    ["Quản trị hệ thống", "Có", "Không", "Có"],
]


# ============================================================
# Ảnh minh hoạ (vẽ bằng PIL - không tải ảnh từ Internet)
# ============================================================

def _ve_so_do_quy_trinh(duong_dan: Path) -> Path:
    """Sơ đồ khối quy trình mượn - dùng để kiểm tra hệ thống có tìm được nội dung nằm
    TRONG hình hay không (chữ trong hình không nằm ở lớp text của tài liệu)."""
    anh = Image.new("RGB", (900, 260), "white")
    ve = ImageDraw.Draw(anh)
    cac_buoc = ["Quét thẻ", "Kiểm tra nợ", "Ghi mượn", "Hẹn ngày trả"]
    x = 20
    for buoc in cac_buoc:
        ve.rectangle([x, 90, x + 180, 170], outline="black", width=3)
        ve.text((x + 30, 122), buoc, fill="black")
        if x + 180 < 800:
            ve.line([x + 180, 130, x + 210, 130], fill="black", width=3)
            ve.line([x + 200, 122, x + 210, 130], fill="black", width=3)
            ve.line([x + 200, 138, x + 210, 130], fill="black", width=3)
        x += 210
    ve.text((20, 20), "So do quy trinh muon tai lieu", fill="black")
    anh.save(duong_dan)
    return duong_dan


def _ve_bieu_do_cot(duong_dan: Path) -> Path:
    """Biểu đồ cột - số liệu CHỈ có trong hình, không lặp lại ở text, để kiểm tra khả năng
    đọc nội dung trực quan (cần vision model; không có nó thì chỉ tìm được qua caption)."""
    anh = Image.new("RGB", (760, 420), "white")
    ve = ImageDraw.Draw(anh)
    du_lieu = [("2021", 120), ("2022", 200), ("2023", 260), ("2024", 330)]
    ve.line([70, 360, 720, 360], fill="black", width=3)
    ve.line([70, 40, 70, 360], fill="black", width=3)
    x = 120
    for nhan, cao in du_lieu:
        ve.rectangle([x, 360 - cao, x + 90, 360], fill="#4a7ebb", outline="black", width=2)
        ve.text((x + 25, 370), nhan, fill="black")
        ve.text((x + 25, 360 - cao - 20), str(cao), fill="black")
        x += 150
    ve.text((70, 12), "Luot muon theo nam (nghin luot)", fill="black")
    anh.save(duong_dan)
    return duong_dan


# ============================================================
# Sinh từng loại tài liệu
# ============================================================

def tao_docx_dai(duong_dan: Path) -> None:
    """Tài liệu DÀI, nhiều cấp tiêu đề - kiểm tra chunking có tôn trọng ranh giới mục không.

    Dùng style "Heading N" chuẩn của Word: đây là tín hiệu cấu trúc ĐÁNG TIN NHẤT trong 3
    định dạng (PDF chỉ suy ra được từ cỡ chữ), nên là mốc so sánh để biết heuristic PDF
    còn cách bao xa so với trường hợp lý tưởng.
    """
    tai_lieu = Document()
    for i, chuong in enumerate(CAC_CHUONG_DAI):
        for noi_dung, cap in chuong:
            if cap == 0:
                tai_lieu.add_paragraph(noi_dung)
            else:
                tai_lieu.add_heading(noi_dung, level=cap)
        # Ngắt trang giữa các chương, KHÔNG ngắt sau chương cuối - ngắt cuối sẽ tạo ra một
        # "trang" rỗng bị loader cảnh báo và bỏ qua (nhiễu log, không có nội dung thật).
        if i < len(CAC_CHUONG_DAI) - 1:
            tai_lieu.paragraphs[-1].add_run().add_break(WD_BREAK.PAGE)
    tai_lieu.save(duong_dan)


def tao_docx_nhieu_muc_tuong_tu(duong_dan: Path) -> None:
    """Tài liệu GÂY NHIỄU: 20 điều khoản gần như trùng nhau, mỗi điều 1 trang.

    Mục đích không phải để đọc, mà để tạo ra tình huống khó thật: khi hỏi về hồ sơ loại
    "G1", có 19 đoạn khác trông y hệt chỉ khác mã và con số. Hệ thống phải bám đúng từ khoá
    hiếm ("G1") thay vì trôi theo độ giống chung chung của cả đoạn.
    """
    tai_lieu = Document()
    tai_lieu.add_heading("Quy chế xử lý hồ sơ", level=1)
    cac_muc = _sinh_muc_tuong_tu()
    for i in range(0, len(cac_muc), 2):
        for noi_dung, cap in cac_muc[i:i + 2]:
            if cap == 0:
                tai_lieu.add_paragraph(noi_dung)
            else:
                tai_lieu.add_heading(noi_dung, level=cap)
        if i + 2 < len(cac_muc):
            tai_lieu.paragraphs[-1].add_run().add_break(WD_BREAK.PAGE)
    tai_lieu.save(duong_dan)


def tao_docx_co_bang(duong_dan: Path, anh_so_do: Path) -> None:
    """Bảng xen giữa văn xuôi + BẢNG LỒNG trong ô + ảnh - gom nhiều bẫy vào 1 file.

    Bảng lồng là giới hạn đã biết của python-docx (document.tables chỉ thấy bảng cấp cao
    nhất). Đưa vào đây có chủ đích: để giới hạn đó hiện ra trong số đo thay vì âm thầm,
    và để sau này nếu có cách xử lý thì có sẵn ca kiểm thử.
    """
    tai_lieu = Document()
    tai_lieu.add_heading("Quy định mượn trả tài liệu", level=1)
    tai_lieu.add_paragraph(
        "Bảng dưới đây quy định thời hạn mượn và mức phí phạt áp dụng cho từng loại tài "
        "liệu. Người đọc cần nắm rõ trước khi làm thủ tục mượn."
    )

    bang = tai_lieu.add_table(rows=len(BANG_PHI_PHAT), cols=len(BANG_PHI_PHAT[0]))
    bang.style = "Table Grid"
    for i, hang in enumerate(BANG_PHI_PHAT):
        for j, o in enumerate(hang):
            bang.cell(i, j).text = o

    tai_lieu.add_heading("Phân quyền người dùng", level=1)
    tai_lieu.add_paragraph(
        "Hệ thống phân biệt ba vai trò với quyền hạn khác nhau trên dữ liệu tài liệu."
    )
    bang2 = tai_lieu.add_table(rows=len(BANG_PHAN_QUYEN), cols=len(BANG_PHAN_QUYEN[0]))
    bang2.style = "Table Grid"
    for i, hang in enumerate(BANG_PHAN_QUYEN):
        for j, o in enumerate(hang):
            bang2.cell(i, j).text = o

    # Bảng LỒNG trong ô của bảng khác - python-docx không liệt kê được qua document.tables.
    o_chua_bang_long = bang2.cell(len(BANG_PHAN_QUYEN) - 1, len(BANG_PHAN_QUYEN[0]) - 1)
    bang_long = o_chua_bang_long.add_table(rows=2, cols=2)
    bang_long.cell(0, 0).text = "Ghi chú"
    bang_long.cell(0, 1).text = "Áp dụng"
    bang_long.cell(1, 0).text = "Cần phê duyệt"
    bang_long.cell(1, 1).text = "Từ 2024"

    tai_lieu.add_heading("Sơ đồ quy trình", level=1)
    tai_lieu.add_paragraph("Hình 1: Sơ đồ các bước xử lý khi người đọc mượn tài liệu.")
    tai_lieu.add_picture(str(anh_so_do), width=DocxInches(5.5))
    tai_lieu.save(duong_dan)


def tao_pptx_ngan(duong_dan: Path) -> None:
    """Tài liệu NGẮN, thưa chữ - cực đối lập với giáo trình dày chữ.

    Mỗi slide chỉ vài dòng nên cả slide thường lọt gọn trong ngân sách ngữ cảnh; đây là
    ca kiểm chứng rằng việc tối ưu cho tài liệu dài KHÔNG làm hỏng tài liệu ngắn.
    """
    trinh_chieu = Presentation()
    bo_cuc = trinh_chieu.slide_layouts[1]  # Title and Content
    for tieu_de, cac_y in NOI_DUNG_SLIDE:
        slide = trinh_chieu.slides.add_slide(bo_cuc)
        slide.shapes.title.text = tieu_de
        khung = slide.placeholders[1].text_frame
        khung.text = cac_y[0]
        for y in cac_y[1:]:
            khung.add_paragraph().text = y
    trinh_chieu.save(duong_dan)


def tao_pptx_bang_anh(duong_dan: Path, anh_so_do: Path, anh_bieu_do: Path) -> None:
    """PPTX có bảng, ảnh thường, và ảnh+bảng NẰM TRONG GROUP SHAPE.

    Group shape là bẫy đã xác minh: vòng lặp `for shape in slide.shapes` phẳng sẽ không
    nhìn thấy gì bên trong group, nên nội dung đó biến mất khỏi index mà không báo lỗi.
    """
    trinh_chieu = Presentation()
    bo_cuc_trong = trinh_chieu.slide_layouts[5]  # Title Only

    slide = trinh_chieu.slides.add_slide(bo_cuc_trong)
    slide.shapes.title.text = "Bảng phí phạt theo loại tài liệu"
    hinh_bang = slide.shapes.add_table(
        len(BANG_PHI_PHAT), len(BANG_PHI_PHAT[0]),
        Inches(0.4), Inches(1.6), Inches(9.0), Inches(3.0),
    )
    for i, hang in enumerate(BANG_PHI_PHAT):
        for j, o in enumerate(hang):
            hinh_bang.table.cell(i, j).text = o

    slide2 = trinh_chieu.slides.add_slide(bo_cuc_trong)
    slide2.shapes.title.text = "Biểu đồ lượt mượn theo năm"
    slide2.shapes.add_picture(str(anh_bieu_do), Inches(1.0), Inches(1.7), width=Inches(6.5))
    khung_caption = slide2.shapes.add_textbox(Inches(1.0), Inches(5.6), Inches(6.5), Inches(0.6))
    khung_caption.text_frame.text = "Hình 2: Lượt mượn tài liệu tăng dần qua các năm."

    # Slide có GROUP SHAPE chứa ảnh + chú thích bên trong.
    slide3 = trinh_chieu.slides.add_slide(bo_cuc_trong)
    slide3.shapes.title.text = "Quy trình mượn (nhóm hình)"
    nhom = slide3.shapes.add_group_shape()
    nhom.shapes.add_picture(str(anh_so_do), Inches(0.6), Inches(1.8), width=Inches(8.0))
    khung_trong_nhom = nhom.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(8.0), Inches(0.6))
    khung_trong_nhom.text_frame.text = (
        "Hình 3: Bốn bước xử lý mượn tài liệu, bắt đầu từ quét thẻ người đọc."
    )
    trinh_chieu.save(duong_dan)


def tao_pdf_hon_hop(duong_dan: Path, anh_bieu_do: Path) -> None:
    """PDF có tiêu đề phân biệt CHỈ bằng cỡ chữ, bảng kẻ khung, và ảnh kèm caption.

    Đây là ca khó nhất: PDF không lưu cấu trúc logic (không có khái niệm "Heading 1" như
    DOCX, không có placeholder title như PPTX) - muốn biết dòng nào là tiêu đề chỉ còn
    cách suy ra từ cỡ chữ và độ dài dòng. File này là dữ liệu để hiệu chỉnh ngưỡng đó.
    """
    pdfmetrics.registerFont(TTFont(_TEN_FONT, str(_FONT_HE_THONG)))
    trang = canvas.Canvas(str(duong_dan), pagesize=A4)
    rong, cao = A4

    def _dong(y, chu, co=11):
        trang.setFont(_TEN_FONT, co)
        trang.drawString(2.5 * cm, y, chu)

    # --- Trang 1: tiêu đề nhiều cấp + văn xuôi ---
    y = cao - 3 * cm
    _dong(y, "Sổ tay nghiệp vụ thư viện", 20); y -= 1.2 * cm
    _dong(y, "Phần 1. Mượn và trả tài liệu", 15); y -= 1.0 * cm
    for cau in [
        "Người đọc phải xuất trình thẻ còn hiệu lực khi mượn tài liệu tại quầy.",
        "Thủ thư kiểm tra số tài liệu đang mượn và các khoản phạt chưa thanh toán.",
        "Mỗi thẻ được mượn tối đa năm tài liệu cùng lúc trong thời hạn hai mươi mốt ngày.",
        "Tài liệu quý hiếm chỉ được đọc tại chỗ và không cho mượn mang về.",
    ]:
        _dong(y, cau); y -= 0.7 * cm
    y -= 0.5 * cm
    _dong(y, "Phần 2. Bảo quản kho sách", 15); y -= 1.0 * cm
    for cau in [
        "Kho sách duy trì nhiệt độ ổn định và độ ẩm trong ngưỡng cho phép.",
        "Ánh sáng trực tiếp làm giấy ố vàng nên khu lưu trữ lâu dài không có cửa sổ.",
        "Kiểm kê định kỳ giúp phát hiện sớm tài liệu xuống cấp cần phục chế.",
    ]:
        _dong(y, cau); y -= 0.7 * cm
    trang.showPage()

    # --- Trang 2: bảng KẺ KHUNG (để pdfplumber.find_tables nhận ra) ---
    y = cao - 3 * cm
    _dong(y, "Phần 3. Biểu phí áp dụng", 15); y -= 1.2 * cm
    trang.setFont(_TEN_FONT, 10)
    rong_cot = [4.5 * cm, 3.2 * cm, 4.2 * cm, 3.6 * cm]
    cao_hang = 0.9 * cm
    y_bang = y
    for i, hang in enumerate(BANG_PHI_PHAT):
        x = 2.0 * cm
        for j, o in enumerate(hang):
            trang.rect(x, y_bang - cao_hang, rong_cot[j], cao_hang)
            trang.drawString(x + 0.15 * cm, y_bang - cao_hang + 0.3 * cm, o)
            x += rong_cot[j]
        y_bang -= cao_hang
    trang.showPage()

    # --- Trang 3: ảnh + caption ---
    y = cao - 3 * cm
    _dong(y, "Phần 4. Thống kê sử dụng", 15); y -= 1.0 * cm
    _dong(y, "Số liệu lượt mượn được tổng hợp cuối mỗi năm để lập kế hoạch bổ sung.")
    trang.drawImage(str(anh_bieu_do), 2.5 * cm, cao - 14 * cm, width=12 * cm,
                    height=6.6 * cm, preserveAspectRatio=True)
    _dong(cao - 15 * cm, "Hình 4: Lượt mượn tài liệu tăng dần qua các năm.", 10)
    trang.showPage()
    trang.save()


def main() -> None:
    config.RAW_DOCS_DIR.mkdir(parents=True, exist_ok=True)
    THU_MUC_ANH_TAM.mkdir(parents=True, exist_ok=True)

    if not _FONT_HE_THONG.exists():
        print(f"Không tìm thấy font {_FONT_HE_THONG} - PDF sẽ mất dấu tiếng Việt.")
        return

    anh_so_do = _ve_so_do_quy_trinh(THU_MUC_ANH_TAM / "so_do_quy_trinh.png")
    anh_bieu_do = _ve_bieu_do_cot(THU_MUC_ANH_TAM / "bieu_do_luot_muon.png")

    cac_viec = [
        ("dai_nhieu_chu.docx", lambda p: tao_docx_dai(p)),
        ("nhieu_muc_tuong_tu.docx", lambda p: tao_docx_nhieu_muc_tuong_tu(p)),
        ("co_bang_anh.docx", lambda p: tao_docx_co_bang(p, anh_so_do)),
        ("ngan_thua_chu.pptx", lambda p: tao_pptx_ngan(p)),
        ("bang_anh_nhom.pptx", lambda p: tao_pptx_bang_anh(p, anh_so_do, anh_bieu_do)),
        ("hon_hop.pdf", lambda p: tao_pdf_hon_hop(p, anh_bieu_do)),
    ]
    for ten, ham_tao in cac_viec:
        duong_dan = config.RAW_DOCS_DIR / ten
        ham_tao(duong_dan)
        print(f"Đã tạo {duong_dan.name} ({duong_dan.stat().st_size // 1024} KB)")

    print(f"\nXong. Tài liệu mẫu nằm ở: {config.RAW_DOCS_DIR}")
    print("Bấm \"Đọc tài liệu\" trên Streamlit (hoặc chạy evaluation) để dùng.")


if __name__ == "__main__":
    main()
