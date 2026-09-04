"""Trích xuất hình ảnh từ tài liệu và gắn chúng với văn bản xung quanh.

Bài toán: nội dung nằm TRONG hình (sơ đồ, biểu đồ, ảnh chụp bảng) hoàn toàn vô hình với
hệ thống chỉ đọc lớp text. Người dùng hỏi "sơ đồ quy trình gồm những bước nào" thì dù tài
liệu có sơ đồ đó, hệ thống vẫn không tìm ra.

Có 2 mức xử lý, tách rời có chủ đích:

  Mức 1 (module này, luôn bật): trích ảnh ra file + lấy VĂN BẢN LÂN CẬN làm chú thích.
      Không cần model nào, gần như miễn phí. Đủ để tìm được hình qua caption kiểu
      "Hình 3: Bốn bước xử lý mượn tài liệu" - vốn là cách tài liệu thật hay đánh số hình.

  Mức 2 (rag/vision_caption.py, tuỳ chọn): cho model vision mô tả nội dung bên trong hình.
      Bắt được cả thứ không có trong caption (số liệu trên biểu đồ, nhãn trong sơ đồ),
      nhưng mỗi ảnh tốn một lượt gọi model nên chỉ bật khi thật sự cần.

Tách 2 mức để tài liệu thuần văn bản không phải trả bất kỳ chi phí nào, còn tài liệu nhiều
hình vẫn dùng được đầy đủ - đúng yêu cầu "hoạt động ổn định trên nhiều loại tài liệu".
"""

import logging
import re
from collections import Counter
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from docx.opc.constants import RELATIONSHIP_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config
from rag.bo_nho_dem import bam_bytes, bam_file

logger = logging.getLogger(__name__)

MOC_ANH = "[HÌNH]"

# Ảnh quá nhỏ gần như luôn là logo, đường kẻ trang trí, bullet icon - không mang nội dung
# nhưng vẫn chiếm 1 vector và 1 lượt gọi model vision nếu bật. Lọc theo kích thước là cách
# rẻ và tổng quát hơn nhiều so với dò tên file hay vị trí.
KICH_THUOC_ANH_TOI_THIEU = 120  # điểm ảnh, áp cho cả chiều rộng lẫn chiều cao

# Số ký tự văn bản lân cận lấy làm chú thích cho ảnh.
DO_DAI_CHU_THICH_LAN_CAN = 400

# Dòng bắt đầu bằng "Hình 3:", "Figure 2.", "Biểu đồ 1 -"... gần như chắc chắn là chú thích
# của hình ngay cạnh. Ưu tiên những dòng này hơn văn bản lân cận chung chung.
_MAU_DONG_CHU_THICH = re.compile(
    r"^\s*(hình|hinh|figure|fig|biểu đồ|bieu do|sơ đồ|so do|bảng|table)\s*\d+\s*[:.\-]",
    re.IGNORECASE,
)


def _ten_file_an_toan(nguon: str, trang: int, thu_tu: int, duoi: str) -> str:
    """Tên file ảnh suy ra từ (nguồn, trang, thứ tự) nên ổn định giữa các lần build - build
    lại index không tạo ra một đống file rác trùng nội dung khác tên."""
    goc = re.sub(r"[^\w\-.]", "_", Path(nguon).stem)[:60]
    return f"{goc}__t{trang}_{thu_tu}{duoi}"


def _chon_chu_thich(cac_dong: List[str]) -> str:
    """Chọn văn bản mô tả hình từ các dòng lân cận.

    Ưu tiên dòng có dạng chú thích chuẩn ("Hình 3: ..."); nếu không có thì gộp văn bản lân
    cận làm ngữ cảnh thay thế. Lý do ưu tiên: dòng chú thích mô tả ĐÚNG hình đó, còn văn bản
    lân cận chỉ là "quanh đây nói về gì" - dùng được nhưng kém chính xác hơn.
    """
    cac_dong = [d.strip() for d in cac_dong if d and d.strip()]
    dong_chu_thich = [d for d in cac_dong if _MAU_DONG_CHU_THICH.match(d)]
    if dong_chu_thich:
        return " ".join(dong_chu_thich)[:DO_DAI_CHU_THICH_LAN_CAN]
    return " ".join(cac_dong)[:DO_DAI_CHU_THICH_LAN_CAN]


def _la_anh_cua_trang_chu(rong: float, cao: float, dien_tich_trang: float) -> bool:
    """Ảnh này có phải là ảnh chụp của một TRANG CHỮ (PDF scan) không?

    Chỉ được gọi cho những trang mà OCR ĐÃ ĐỌC RA CẢ MỘT TRANG CHỮ (xem doc_pdf). Nói cách
    khác, câu hỏi "trang này là chữ hay là hình" KHÔNG được trả lời bằng cách đoán qua kích
    thước ảnh, mà bằng kết quả đo: OCR lấy ra được bao nhiêu chữ từ chính trang đó.

    Vì sao phải làm vậy - hai lần thử trước đều sai theo hai kiểu khác nhau:
      - Chỉ so DIỆN TÍCH với ngưỡng 0.85: trượt ngay, vì ảnh scan của một giáo trình thật
        chỉ phủ 81.2% trang (sách có lề ~50pt mỗi bên).
      - Hạ ngưỡng xuống 0.6 kèm điều kiện "trang không có chữ": bắt được sách scan, nhưng
        nuốt luôn 5 ảnh nền của slide tiêu đề và bìa sách trong corpus cũ - những trang có
        rất ít chữ nhưng ảnh của chúng KHÔNG phải là ảnh chụp chữ.

    Điều phân biệt hai trường hợp đó không phải kích thước, mà là NỘI DUNG của ảnh: ảnh chụp
    trang sách chứa cả nghìn ký tự chữ, còn ảnh nền của slide tiêu đề thì không. Chỉ OCR mới
    biết được điều đó, nên quyết định phải đi sau OCR.

    Phần diện tích còn lại chỉ là chốt phụ: trên một trang scan vẫn có thể có thêm một hình
    nhỏ được nhúng riêng, và hình nhỏ đó thì vẫn nên giữ.
    """
    if dien_tich_trang <= 0:
        return False
    return (rong * cao) / dien_tich_trang >= config.TY_LE_DIEN_TICH_ANH_TOAN_TRANG


def ly_do_loai_anh(rong: float, cao: float, dien_tich_trang: float = 0.0) -> Optional[str]:
    """Ảnh này có đáng đưa vào index không? Trả về LÝ DO loại, hoặc None nếu giữ.

    Ba chốt, xếp theo thứ tự từ rẻ tới đắt và từ chắc chắn tới phỏng đoán:

      1. KÍCH THƯỚC TUYỆT ĐỐI (KICH_THUOC_ANH_TOI_THIEU) - chốt đã có từ trước.
      2. TỈ LỆ CẠNH (config.TY_LE_CANH_ANH_TRANG_TRI) - một hình dẹt 20:1 là đường kẻ hoặc
         thanh màu trang trí, không phải sơ đồ hay biểu đồ. Chốt này bắt được đúng loại ảnh
         mà chốt kích thước bỏ lọt: một đường kẻ ngang dưới tiêu đề slide rộng 900px nên
         vượt xa ngưỡng 120px, nhưng cao có 6px.
      3. TỈ LỆ DIỆN TÍCH SO VỚI TRANG (config.TY_LE_DIEN_TICH_ANH_TOI_THIEU) - chỉ áp dụng
         khi biết diện tích trang (PDF). Icon và logo góc trang chiếm rất ít diện tích;
         hình mang nội dung thật thì gần như luôn được in đủ to để người đọc nhìn được.

    Vì sao lọc TRƯỚC khi render và trước khi gọi model vision, chứ không lọc ở cuối luồng:
    mỗi ảnh được giữ lại kéo theo một lượt render, một file trên đĩa, một lượt gọi model
    vision (~1,9 giây theo benchmark của project) và một vector trong index. Loại một icon
    ở cuối luồng thì bốn khoản chi phí đó đã tiêu mất rồi.
    """
    if rong < KICH_THUOC_ANH_TOI_THIEU or cao < KICH_THUOC_ANH_TOI_THIEU:
        return "quá nhỏ"
    canh_dai, canh_ngan = max(rong, cao), max(min(rong, cao), 1e-9)
    if config.TY_LE_CANH_ANH_TRANG_TRI > 0 and canh_dai / canh_ngan > config.TY_LE_CANH_ANH_TRANG_TRI:
        return "dải trang trí (tỉ lệ cạnh quá dẹt)"
    if dien_tich_trang > 0 and config.TY_LE_DIEN_TICH_ANH_TOI_THIEU > 0:
        if (rong * cao) / dien_tich_trang < config.TY_LE_DIEN_TICH_ANH_TOI_THIEU:
            return "chiếm quá ít diện tích trang (icon/logo)"
    return None


def ly_do_loai_anh_blob(du_lieu: bytes) -> Optional[str]:
    """Như ly_do_loai_anh() nhưng cho ảnh PPTX/DOCX - đọc kích thước thẳng từ bytes.

    PPTX/DOCX cho ảnh dưới dạng blob nhúng sẵn, không có khái niệm "diện tích trang" để so,
    nên chỉ áp được hai chốt hình dạng. Đọc kích thước bằng PIL chỉ tốn phần HEADER của ảnh
    (vài chục byte đầu), không giải nén cả bức - rẻ hơn nhiều so với việc ghi file ra đĩa
    rồi mới phát hiện đó là một cái logo.

    Trả None (tức GIỮ ảnh) khi không đọc được kích thước: định dạng lạ như SVG hay EMF không
    mở được bằng PIL, và bỏ một hình thật vì lý do "không đo được" thì tệ hơn hẳn giữ thừa
    một logo - bước loc_anh_lap_lai() phía sau vẫn còn cơ hội bắt nó.
    """
    try:
        from io import BytesIO

        from PIL import Image

        with Image.open(BytesIO(du_lieu)) as anh:
            rong, cao = anh.size
    except Exception:  # noqa: BLE001 - định dạng lạ / ảnh hỏng thì cứ giữ
        return None
    return ly_do_loai_anh(float(rong), float(cao))


def loc_anh_lap_lai(cac_ban_ghi: List[Dict], nguon: str) -> List[Dict]:
    """Loại những ảnh có NỘI DUNG GIỐNG HỆT lặp lại nhiều lần trong cùng một tài liệu.

    Đây là logo trường, watermark, khung viền mẫu slide - thứ xuất hiện trên mọi trang và
    tạo ra hàng chục chunk giống hệt nhau trong index. Chúng lọt qua mọi chốt hình dạng ở
    ly_do_loai_anh() vì bản thân chúng là ảnh to, rõ, tỉ lệ cạnh bình thường; thứ duy nhất
    tố cáo chúng là VIỆC LẶP LẠI.

    So sánh bằng băm nội dung file (không phải tên file, không phải kích thước): hai bản
    sao của cùng một logo được trích ra ở hai trang khác nhau sẽ có tên khác nhau nhưng
    byte giống hệt nhau. Xem config.SO_LAN_LAP_COI_LA_LOGO để biết vì sao ngưỡng là 4.

    Bản ghi được gắn thêm "bam_anh" để bước chú thích vision sau đó dùng lại làm khoá cache
    mà không phải băm lại file lần nữa.
    """
    if not cac_ban_ghi or config.SO_LAN_LAP_COI_LA_LOGO <= 0:
        return cac_ban_ghi

    for ban_ghi in cac_ban_ghi:
        if not ban_ghi.get("bam_anh"):
            try:
                ban_ghi["bam_anh"] = bam_file(Path(ban_ghi["duong_dan_anh"]))
            except OSError:
                ban_ghi["bam_anh"] = ""

    so_lan = Counter(b["bam_anh"] for b in cac_ban_ghi if b["bam_anh"])
    bam_logo = {b for b, n in so_lan.items() if n >= config.SO_LAN_LAP_COI_LA_LOGO}
    if not bam_logo:
        return cac_ban_ghi

    giu = [b for b in cac_ban_ghi if b["bam_anh"] not in bam_logo]
    logger.info(
        "'%s': bỏ %d bản ghi ảnh thuộc %d hình lặp lại từ %d lần trở lên (logo/watermark/"
        "khung mẫu slide) - chúng chỉ tạo ra các chunk giống hệt nhau trong index.",
        nguon, len(cac_ban_ghi) - len(giu), len(bam_logo), config.SO_LAN_LAP_COI_LA_LOGO,
    )
    return giu


def _ban_ghi_anh(
    nguon: str, trang: int, duong_dan_anh: Path, chu_thich: str, bam_anh: str = ""
) -> Dict:
    """Một ảnh trở thành một "trang" riêng trong luồng dữ liệu.

    Dùng lại đúng schema {nguon, trang, noidung} của document_loader thay vì tạo luồng
    riêng: nhờ vậy ảnh đi qua chunking/embedding/vector_store/citation y hệt văn bản, không
    module nào phải biết đến khái niệm "ảnh". Chỉ 2 trường phụ được thêm để UI hiển thị
    được ảnh và để chunking biết không cắt nhỏ bản ghi này.

    Bản ghi ở đây vẫn được tạo kể cả khi chưa có chú thích nào: chú thích bằng model vision
    được thêm SAU, ở doc_thu_muc(), nên bỏ ảnh ngay tại đây sẽ giết luôn những ảnh sắp được
    model đọc nội dung. Việc loại bỏ ảnh rốt cuộc vẫn rỗng nằm ở
    document_loader._bo_ban_ghi_anh_rong(), chạy sau bước chú thích.
    """
    return {
        "nguon": nguon,
        "trang": trang,
        "noidung": f"{MOC_ANH} {chu_thich}".strip(),
        "loai_noi_dung": "anh",
        "duong_dan_anh": str(duong_dan_anh),
        # Băm NỘI DUNG ảnh - khoá dùng chung cho cả việc phát hiện hình lặp lại
        # (loc_anh_lap_lai) lẫn cache chú thích vision. Với PPTX/DOCX thì byte của ảnh đã
        # nằm sẵn trong RAM lúc này nên băm ngay tại đây là miễn phí; đọc lại từ đĩa sau đó
        # chỉ để băm sẽ tốn thêm một lượt I/O cho mỗi ảnh, mà một bài giảng có thể có tới
        # vài trăm ảnh. Rỗng thì loc_anh_lap_lai() tự băm từ file (đường của PDF, nơi ảnh
        # được render chứ không có sẵn blob).
        "bam_anh": bam_anh,
    }


def ung_vien_anh_trang(trang) -> List[Tuple[Tuple[float, float, float, float], bool]]:
    """Liệt kê ảnh ĐÁNG GIỮ trên một trang PDF - CHƯA render gì cả.

    Trả về [(bbox, la_anh_phu_ca_trang)]. Cờ thứ hai để chỗ gọi tự quyết định sau: một ảnh
    phủ gần kín trang CHỈ nên bị loại khi OCR đã chứng minh trang đó là ảnh chụp một trang
    CHỮ (xem _la_anh_cua_trang_chu) - mà điều đó thì phải chờ OCR chạy xong mới biết. Trả cờ
    ra ngoài cho phép giữ nguyên đúng quy tắc cũ mà vẫn chỉ duyệt PDF một lượt duy nhất.

    Tách "chọn ảnh nào" khỏi "render ảnh đó" là điều làm nên luồng đọc một-lượt: bước chọn
    chỉ đọc metadata đã có sẵn trong đối tượng trang (toạ độ, kích thước), rẻ tới mức chạy
    được ngay trong vòng lặp đọc text; còn bước render mới là phần đắt, nên chỉ chạy cho
    những ảnh đã qua được mọi bộ lọc. Bản trước làm ngược lại - duyệt lại toàn bộ PDF một
    lượt nữa và render trước, lọc sau.
    """
    dien_tich_trang = float(trang.width or 0) * float(trang.height or 0)
    ket_qua = []
    for anh in trang.images:
        rong = float(anh.get("width") or 0)
        cao = float(anh.get("height") or 0)
        ly_do = ly_do_loai_anh(rong, cao, dien_tich_trang)
        if ly_do:
            continue
        # Giới hạn bbox trong khung trang: ảnh tràn mép trang sẽ khiến crop() ném lỗi.
        bbox = (
            max(anh["x0"], trang.bbox[0]), max(anh["top"], trang.bbox[1]),
            min(anh["x1"], trang.bbox[2]), min(anh["bottom"], trang.bbox[3]),
        )
        ket_qua.append((bbox, _la_anh_cua_trang_chu(rong, cao, dien_tich_trang)))
    return ket_qua


def luu_anh_trang_pdf(
    nguon: str, trang, so_trang: int, cac_bbox: List[Tuple[float, float, float, float]],
    cac_dong_text: List[str],
) -> List[Dict]:
    """Render + lưu ra file những ảnh đã được ung_vien_anh_trang() chọn.

    pdfplumber không giải mã được stream ảnh gốc, nên phải RENDER LẠI vùng chứa ảnh thành
    pixel (crop -> to_image). Đánh đổi: không giữ được file gốc nguyên vẹn, nhưng đủ tốt để
    model vision đọc và để hiển thị cho người dùng đối chiếu - đúng 2 mục đích đang cần.
    """
    ket_qua = []
    chu_thich = _chon_chu_thich(cac_dong_text)
    for thu_tu, bbox in enumerate(cac_bbox, start=1):
        try:
            pil = trang.crop(bbox).to_image(resolution=110).original
        except Exception as loi:  # noqa: BLE001 - pdfplumber ném nhiều loại lỗi khác nhau
            logger.warning("Bỏ qua 1 ảnh ở trang %d của '%s': %s",
                           so_trang, nguon, type(loi).__name__)
            continue
        dich = config.IMAGES_DIR / _ten_file_an_toan(nguon, so_trang, thu_tu, ".png")
        pil.save(dich)
        ket_qua.append(_ban_ghi_anh(nguon, so_trang, dich, chu_thich))
    return ket_qua


def trich_anh_pptx(duong_dan: Path, trinh_chieu) -> List[Dict]:
    """Trích ảnh từng slide, kể cả ảnh nằm trong group shape."""
    from rag.document_loader import duyet_shape  # nhập tại chỗ để tránh phụ thuộc vòng

    ket_qua = []
    for so_slide, slide in enumerate(trinh_chieu.slides, start=1):
        cac_shape = list(duyet_shape(slide.shapes))
        cac_dong = [
            s.text_frame.text for s in cac_shape
            if s.has_text_frame and s.text_frame.text.strip()
        ]
        thu_tu = 0
        da_lay = set()  # partname của ảnh đã lấy, để bước dự phòng không lấy trùng
        for shape in cac_shape:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            # shape.image ném lỗi trong 2 trường hợp thực tế đã gặp: ảnh LIÊN KẾT ngoài
            # (Insert > Link to File, không có byte nào trong file), và ảnh SVG do
            # PowerPoint đời mới chèn (blip trỏ tới ảnh qua phần mở rộng mà python-pptx
            # không lần ra được). Không bắt lỗi ở đây thì cả lần build index hỏng.
            try:
                anh = shape.image
            except (ValueError, KeyError, AttributeError):
                continue  # để bước quét rels bên dưới lo
            if ly_do_loai_anh_blob(anh.blob):
                # Đánh dấu là ĐÃ XỬ LÝ để bước quét rels dự phòng không nhặt lại đúng cái
                # logo vừa bị loại - hai đường trích ảnh phải thống nhất với nhau về việc
                # ảnh nào đáng giữ, nếu không bộ lọc chỉ có tác dụng ở một nửa số trường hợp.
                da_lay.add(getattr(anh, "sha1", None) or "")
                continue
            thu_tu += 1
            dich = config.IMAGES_DIR / _ten_file_an_toan(
                duong_dan.name, so_slide, thu_tu, f".{anh.ext}"
            )
            dich.write_bytes(anh.blob)
            da_lay.add(getattr(anh, "sha1", None) or dich.name)
            ket_qua.append(
                _ban_ghi_anh(
                    duong_dan.name, so_slide, dich, _chon_chu_thich(cac_dong),
                    bam_bytes(anh.blob),
                )
            )

        # DỰ PHÒNG: quét quan hệ (rels) của slide để bắt những ảnh mà API shape không với
        # tới được. Đo trên một bài giảng thật: 11 shape ảnh đều ném lỗi qua API shape,
        # nhưng gói .pptx chứa 8 ảnh nhúng thật (9 MB, gồm cả SVG) - nếu chỉ dựa vào API
        # shape thì toàn bộ số ảnh đó biến mất khỏi index mà không có dấu hiệu gì.
        # Cùng kỹ thuật đã dùng cho DOCX, nhưng ở đây quét theo TỪNG SLIDE nên vẫn giữ
        # được thông tin ảnh thuộc slide nào.
        for quan_he in slide.part.rels.values():
            if quan_he.reltype != RELATIONSHIP_TYPE.IMAGE or quan_he.is_external:
                continue
            try:
                phan_anh = quan_he.target_part
                khoa = getattr(phan_anh, "sha1", None) or str(phan_anh.partname)
                if khoa in da_lay:
                    continue
                du_lieu = phan_anh.blob
                duoi = Path(str(phan_anh.partname)).suffix or ".png"
            except Exception as loi:  # noqa: BLE001 - gói hỏng thì bỏ qua đúng ảnh đó
                logger.warning(
                    "Slide %d của '%s': bỏ qua 1 ảnh (%s).",
                    so_slide, duong_dan.name, type(loi).__name__,
                )
                continue
            da_lay.add(khoa)
            if ly_do_loai_anh_blob(du_lieu):
                continue
            thu_tu += 1
            dich = config.IMAGES_DIR / _ten_file_an_toan(
                duong_dan.name, so_slide, thu_tu, duoi
            )
            dich.write_bytes(du_lieu)
            ket_qua.append(
                _ban_ghi_anh(
                    duong_dan.name, so_slide, dich, _chon_chu_thich(cac_dong),
                    bam_bytes(du_lieu),
                )
            )
    return ket_qua


def trich_anh_docx(duong_dan: Path, document) -> List[Dict]:
    """Trích ảnh DOCX qua quan hệ (rels) của phần thân tài liệu.

    Không dùng `document.inline_shapes` vì bộ đó CHỈ chứa ảnh dạng inline - ảnh neo/nổi
    (floating, kiểu bọc chữ quanh ảnh) không nằm trong đó và sẽ bị bỏ sót. Duyệt rels bắt
    được cả hai loại.

    Đánh đổi: rels không cho biết ảnh nằm ở trang nào, nên mọi ảnh DOCX được gán trang 1.
    Chấp nhận được vì bản thân "trang" trong DOCX cũng chỉ là ước lượng theo ngắt trang
    cứng (xem doc_docx), và chú thích lân cận mới là thứ giúp tìm ra hình.
    """
    ket_qua = []
    cac_dong = [p.text for p in document.paragraphs if p.text.strip()]
    thu_tu = 0
    for quan_he in document.part.rels.values():
        if quan_he.reltype != RELATIONSHIP_TYPE.IMAGE:
            continue
        try:
            du_lieu = quan_he.target_part.blob
            duoi = f".{quan_he.target_part.image.ext}"
        except Exception as loi:  # noqa: BLE001
            logger.warning("Bỏ qua 1 ảnh trong '%s': %s", duong_dan.name, type(loi).__name__)
            continue
        if ly_do_loai_anh_blob(du_lieu):
            continue
        thu_tu += 1
        dich = config.IMAGES_DIR / _ten_file_an_toan(duong_dan.name, 1, thu_tu, duoi)
        dich.write_bytes(du_lieu)
        ket_qua.append(
            _ban_ghi_anh(
                duong_dan.name, 1, dich, _chon_chu_thich(cac_dong), bam_bytes(du_lieu)
            )
        )
    return ket_qua
