"""Đọc tài liệu PDF/PPTX/DOCX, giữ metadata (tên file, số trang/slide) ngay từ bước đọc.

Mỗi trang PDF, slide PPTX, hoặc "trang" DOCX (xem doc_docx) được trả về dưới dạng 1 dict:
    {"nguon": <tên file>, "trang": <số trang/slide, bắt đầu từ 1>, "noidung": <text>}

Trang/slide không có text (ví dụ ảnh scan) bị bỏ qua và chỉ log cảnh báo -
đồ án không yêu cầu OCR nên không thể trích xuất nội dung từ ảnh.
"""

import logging
import re
import unicodedata
from collections import Counter
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Dict, Iterator, List, Optional

import pdfplumber
from docx import Document
from docx.table import Table as BangDocx
from docx.text.paragraph import Paragraph as DoanVanDocx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config
from rag import bo_nho_dem, do_thoi_gian, tai_nguyen_gpu
from rag.image_extractor import (
    MOC_ANH,
    loc_anh_lap_lai,
    luu_anh_trang_pdf,
    trich_anh_docx,
    trich_anh_pptx,
    ung_vien_anh_trang,
)
from rag.vision_caption import (
    bo_sung_chu_thich_vision,
    mo_hinh_vision_co_san,
    ocr_trang_pdf,
    trang_can_ocr,
)

logger = logging.getLogger(__name__)

CAC_DUOI_HO_TRO = (".pdf", ".pptx", ".docx")

# Client Ollama dùng chung cho OCR dự phòng - tạo một lần cho cả lần build thay vì mỗi
# trang một lần. _da_canh_bao_vision để chỉ cảnh báo "chưa pull model" đúng một lần, thay
# vì lặp lại trên từng trang hỏng (sách toán có thể có hàng trăm trang như vậy).
_client_vision = None
_da_canh_bao_vision = False

# Mốc đánh dấu khối bảng trong chuỗi noidung. Có 2 mục đích:
#   1. chunking.py cắt theo mốc này để giữ NGUYÊN VẸN một bảng trong 1 chunk - bảng bị cắt
#      ngang giữa các hàng thì mất luôn quan hệ hàng-cột, tức mất đúng thứ khiến nó là bảng.
#   2. LLM nhìn thấy mốc nên biết đoạn đó là bảng chứ không phải văn xuôi lộn xộn.
# Cố ý KHÔNG dùng dạng số trong ngoặc vuông ([1], [2]) vì trùng cú pháp trích dẫn mà
# citation.py đang dò (_MAU_THAM_CHIEU) - sẽ khiến trích dẫn bắt nhầm.
MOC_BANG_MO = "[BẢNG]"
MOC_BANG_DONG = "[/BẢNG]"

# Watermark nhúng bởi StuDocu (studocu.com) khi tài liệu được chia sẻ dạng "xem trước" -
# lặp lại trên gần như mọi trang của các PDF tải từ nguồn này (dạng "lOMoARcPSD|<số>" ở đầu
# trang, "Downloaded by <tên> (<email>)" ở cuối trang), không mang thông tin thật, chỉ gây
# nhiễu cho embedding (mọi chunk đều dính chung 1 đoạn text giống hệt nhau, làm giảm độ
# phân biệt ngữ nghĩa giữa các chunk).
_MAU_WATERMARK_STUDOCU = re.compile(r"lOMoARcPSD\|\d+|Downloaded by .+?\([^)]*@[^)]*\)")

# Mã ký tự PDF chưa giải mã được, dạng "(cid:10)" - xem giải thích ở _don_dep_watermark().
_MAU_CID_PDF = re.compile(r"\(cid:\d+\)")

# Cụm chữ cái liền không có khoảng trắng - dấu hiệu pdfplumber đã nuốt mất khoảng trắng
# giữa các từ (xem config.BAT_DOC_LAI_TRANG_DINH_CHU). Chỉ đếm chữ cái Latin + chữ có dấu
# tiếng Việt; chữ số và ký hiệu toán không tính vì công thức vốn dĩ viết liền là đúng.
_MAU_CUM_CHU = re.compile(r"[A-Za-zÀ-ỹ]+")


def _ty_le_dinh_chu(text: str) -> float:
    """Tỉ lệ ký tự chữ nằm trong những cụm dài bất thường (không có khoảng trắng ngăn từ).

    Đây là cách đo "trang này có bị dính chữ không" mà KHÔNG cần từ điển: một trang tiếng
    Anh hay tiếng Việt bình thường gần như không có cụm chữ cái dài như vậy, còn trang bị
    nuốt khoảng trắng thì gần một nửa số ký tự nằm trong các cụm như thế.

    HAI CHỐT CHỐNG BÁO ĐỘNG GIẢ, cả hai đều đến từ ca hỏng đo được chứ không phải phòng xa:

    1. Độ dài cụm (config.DO_DAI_CUM_DINH_CHU = 20). Ở mức 15, từ ghép kỹ thuật hợp lệ như
       "Backpropagation" đã bị tính là dính. Nâng lên 20 thì nó ra 0% trong khi trang Bishop
       vẫn 41.7% - khoảng cách giữa hai nhóm còn rộng hơn trước.

    2. Lượng chữ tối thiểu (config.SO_KY_TU_TOI_THIEU_DE_DO). Đây mới là chốt quan trọng:
       nâng độ dài cụm KHÔNG đủ, vì một câu ngắn chứa đúng một từ dài hợp lệ
       ("internationalization of counterrevolutionaries") vẫn cho tới 81.6% ở mọi ngưỡng độ
       dài - đơn giản vì mẫu quá nhỏ. Tỉ lệ chỉ có nghĩa khi có đủ chữ để mà tính tỉ lệ.
       Dưới ngưỡng thì trả 0.0: một trang ít chữ đến thế cũng không đáng đọc lại.

    Đo lại trên corpus thật với cấu hình này: Bishop 41.7%, tám file PDF còn lại 0.0-1.5%.
    """
    do_dai = [len(t) for t in _MAU_CUM_CHU.findall(text)]
    tong = sum(do_dai)
    if tong < config.SO_KY_TU_TOI_THIEU_DE_DO:
        return 0.0
    return sum(d for d in do_dai if d >= config.DO_DAI_CUM_DINH_CHU) / tong


def _ty_le_tu_le(text: str) -> float:
    """Tỉ lệ "từ" chỉ có ĐÚNG MỘT chữ cái - dấu hiệu ngược lại của dính chữ.

    Khi x_tolerance nhỏ hơn khoảng cách chữ thật, pdfplumber chèn khoảng trắng vào GIỮA từ
    ("befor e", "c l a s s"), làm văn bản vỡ vụn thành các chữ cái rời. Đây là cái giá phải
    trả khi ép tham số xuống quá thấp, và là thứ phải đo được thì mới biết mình đang chữa
    bệnh hay đang gây bệnh.
    """
    cac_tu = _MAU_CUM_CHU.findall(text)
    if not cac_tu:
        return 0.0
    return sum(1 for t in cac_tu if len(t) == 1) / len(cac_tu)


def _trich_text(doi_tuong, x_tolerance=None) -> str:
    """extract_text() có/không tham số x_tolerance - gom vào 1 chỗ để 2 nhánh dùng chung."""
    if x_tolerance is None:
        return doi_tuong.extract_text() or ""
    return doi_tuong.extract_text(x_tolerance=x_tolerance) or ""


class HieuChinhXTolerance:
    """Nhớ x_tolerance đã dò được cho MỘT tài liệu, để các trang sau không phải dò lại.

    VÌ SAO CÓ LỚP NÀY. Bản trước dò lại từ đầu trên MỌI trang bị dính chữ: thử lần lượt cả
    4 mức trong CAC_X_TOLERANCE_THU, mỗi mức là một lượt extract_text() đọc lại toàn bộ
    trang, rồi mới chọn mức tốt nhất. Với giáo trình Bishop - nơi gần như trang nào cũng
    dính chữ vì font Computer Modern - đó là gấp 5 lần chi phí đọc text cho cả cuốn sách,
    để rồi lần nào cũng chọn ra đúng một giá trị.

    Điều làm cho việc nhớ lại là HỢP LỆ chứ không phải mẹo tiết kiệm: nguyên nhân dính chữ
    là FONT VÀ CỠ CHỮ của tài liệu (xem _trich_text_thich_ung), mà hai thứ đó gần như không
    đổi trong cùng một cuốn sách. Giá trị tốt cho trang 30 gần như chắc chắn cũng tốt cho
    trang 31.

    Điều làm cho nó AN TOÀN: giá trị đã nhớ chỉ được thử TRƯỚC. Nó phải vượt qua đúng hai
    phép kiểm tra như mọi ứng viên khác (giảm được độ dính, không làm vỡ từ thêm) thì mới
    được dùng; không đạt thì trang đó vẫn dò đầy đủ như cũ. Nhờ vậy một tài liệu trộn nhiều
    font - phụ lục scan, chương chèn từ nguồn khác - không bị đọc hỏng.
    """

    def __init__(self):
        self.x_da_chon: Optional[float] = None
        self.so_lan_dong_y = 0

    @property
    def da_hieu_chinh(self) -> bool:
        return (
            self.x_da_chon is not None
            and self.so_lan_dong_y >= config.SO_TRANG_HIEU_CHINH_X_TOLERANCE
        )

    def thu_tu_uu_tien(self) -> List[float]:
        """Thứ tự thử các mức x_tolerance cho trang tiếp theo."""
        if self.x_da_chon is None:
            return list(config.CAC_X_TOLERANCE_THU)
        con_lai = [x for x in config.CAC_X_TOLERANCE_THU if x != self.x_da_chon]
        return [self.x_da_chon] + con_lai

    def ghi_nhan(self, x_tolerance: float) -> None:
        if x_tolerance == self.x_da_chon:
            self.so_lan_dong_y += 1
        else:
            self.x_da_chon, self.so_lan_dong_y = x_tolerance, 1


def _trich_text_thich_ung(
    doi_tuong, ten_file: str = "", so_trang=None, hieu_chinh: Optional[HieuChinhXTolerance] = None
) -> str:
    """Đọc text một trang PDF, TỰ DÒ tham số đọc lại khi phát hiện chữ bị dính liền.

    Bài toán: pdfplumber quyết định "có khoảng trắng giữa 2 ký tự không" bằng x_tolerance
    (mặc định 3 điểm). Font đặt sát nhau (Computer Modern của sách LaTeX là ca điển hình)
    có khoảng cách từ nhỏ hơn 3 điểm nên mọi khoảng trắng bị nuốt. Nhưng hạ x_tolerance quá
    tay lại gây đúng lỗi ngược: khoảng trắng chèn vào GIỮA từ.

    Vì sao DÒ chứ không dùng một hằng số: giá trị tốt nhất phụ thuộc vào FONT VÀ CỠ CHỮ của
    từng tài liệu, không phải một con số đúng cho mọi PDF. Chọn cứng một giá trị đo được
    trên một cuốn sách là tự buộc mình vào đúng cuốn sách đó - tài liệu khác có mật độ chữ
    khác sẽ cần giá trị khác, và không ai biết để chỉnh.

    Cách làm, áp dụng cho TỪNG TRANG một cách độc lập:
      1. Đọc bình thường. Không dính (dưới ngưỡng) -> xong, không đụng gì tới nó.
      2. Dính -> thử lần lượt các mức trong config.CAC_X_TOLERANCE_THU, ít can thiệp trước.
      3. Chỉ chấp nhận một mức khi nó vừa GIẢM ĐƯỢC độ dính vừa KHÔNG làm vỡ từ thêm quá
         config.MUC_TANG_TU_LE_CHAP_NHAN so với bản gốc. Điều kiện thứ hai chính là cái
         chốt an toàn: nó biến việc "tôi đã tự tay kiểm tra là không hại file khác" thành
         một phép kiểm tra chạy trên mọi trang của mọi tài liệu.
      4. DỪNG SỚM ngay khi một mức đưa độ dính xuống dưới ngưỡng - không thử tiếp các mức
         còn lại. Bản trước luôn thử hết cả 4 mức rồi mới chọn mức giảm dính NHIỀU NHẤT;
         chênh lệch giữa "đạt ngưỡng" và "tốt nhất tuyệt đối" không đổi được nội dung đọc
         ra một cách có ý nghĩa, trong khi cái giá là đọc lại toàn bộ trang thêm 3 lần nữa.
      5. HIỆU CHỈNH THEO TÀI LIỆU: mức đã dùng được ở các trang trước được thử TRƯỚC TIÊN
         (xem HieuChinhXTolerance). Với tài liệu dính chữ toàn tập, từ trang thứ tư trở đi
         hầu như chỉ còn đúng 1 lượt đọc lại thay vì 4.

    Không mức nào đạt thì giữ nguyên bản gốc.
    """
    text = _trich_text(doi_tuong)
    if not config.BAT_DOC_LAI_TRANG_DINH_CHU:
        return text
    ty_le_goc = _ty_le_dinh_chu(text)
    if ty_le_goc < config.TY_LE_DINH_CHU_DE_DOC_LAI:
        return text

    tran_tu_le = _ty_le_tu_le(text) + config.MUC_TANG_TU_LE_CHAP_NHAN
    thu_tu = hieu_chinh.thu_tu_uu_tien() if hieu_chinh else list(config.CAC_X_TOLERANCE_THU)
    tot_nhat, ty_le_tot_nhat, x_tot_nhat = None, ty_le_goc, None
    for x_tolerance in thu_tu:
        thu = _trich_text(doi_tuong, x_tolerance=x_tolerance)
        if not thu:
            continue
        ty_le_thu = _ty_le_dinh_chu(thu)
        if ty_le_thu >= ty_le_tot_nhat or _ty_le_tu_le(thu) > tran_tu_le:
            continue
        tot_nhat, ty_le_tot_nhat, x_tot_nhat = thu, ty_le_thu, x_tolerance
        if ty_le_thu <= config.TY_LE_DINH_CHU_DAT_YEU_CAU:
            # ĐÃ SẠCH - thử tiếp chỉ tốn thêm lượt đọc trang mà không cứu thêm được chữ nào.
            # Ngưỡng ở đây chặt hơn hẳn ngưỡng kích hoạt đọc lại, và điều đó là bắt buộc:
            # xem config.TY_LE_DINH_CHU_DAT_YEU_CAU để biết lỗi đã gặp khi dùng chung một số.
            break

    if tot_nhat is None:
        # Không mức nào vừa đỡ dính vừa không làm vỡ từ -> trang này thật sự viết liền, hoặc
        # font hỏng theo kiểu khác. Giữ bản gốc, đừng đánh đổi lấy một bản cũng hỏng.
        return text
    if hieu_chinh is not None:
        hieu_chinh.ghi_nhan(x_tot_nhat)
    # Sau khi đã hiệu chỉnh xong, mỗi trang dính chữ là chuyện bình thường của tài liệu này
    # chứ không còn là phát hiện đáng báo - hạ xuống DEBUG để log không bị một cuốn sách 700
    # trang nhấn chìm (đây chính là thứ khiến log của lần build cũ không đọc nổi).
    ghi_log = logger.debug if (hieu_chinh and hieu_chinh.da_hieu_chinh) else logger.info
    ghi_log(
        "Trang %s của '%s' bị dính chữ (%.0f%%) - đọc lại với x_tolerance=%.1f, còn %.0f%%.",
        so_trang if so_trang is not None else "?", ten_file or "?",
        ty_le_goc * 100, x_tot_nhat, ty_le_tot_nhat * 100,
    )
    return tot_nhat
# Gộp khoảng trắng liên tiếp còn lại sau khi xoá, nhưng GIỮ NGUYÊN dấu xuống dòng: xuống
# dòng là ranh giới mà chunking dựa vào để cắt (CHUNK_SEPARATORS), xoá đi sẽ làm cả trang
# dính thành một khối và splitter mất chỗ cắt hợp lý.
_MAU_KHOANG_TRANG_THUA = re.compile(r"[ \t]{2,}")

# Vài cụm từ mở đầu trang Mục lục phổ biến (VI/EN) - chỉ giúp nhận diện trang ĐẦU TIÊN của
# mục lục (trang có tiêu đề). Mục lục chi tiết (liệt kê tới từng mục nhỏ) thường dài NHIỀU
# trang - các trang tiếp theo không lặp lại tiêu đề nên cần thêm cách nhận diện khác, xem
# _MAU_DONG_KET_THUC_BANG_SO bên dưới.
_CAC_TIEU_DE_MUC_LUC = ("mục lục", "table of contents", "contents")

# Mỗi dòng mục lục có dạng "<tên mục>... <số trang>" - dùng tỉ lệ dòng kết thúc bằng 1 con
# số để nhận diện các trang mục lục TIẾP THEO (không có tiêu đề "Mục lục" lặp lại). Ngưỡng
# 0.5 đã hiệu chỉnh bằng dữ liệu thật: trang mục lục thật cho tỉ lệ 0.59-0.93, trang nội
# dung thật (văn xuôi, điều luật, danh sách tác giả...) chỉ 0.0-0.07 - cách biệt rất lớn.
_MAU_DONG_KET_THUC_BANG_SO = re.compile(r"\d{1,4}\s*$")
_NGUONG_TY_LE_MUC_LUC = 0.5
_SO_DONG_TOI_THIEU_DE_XET = 6  # trang quá ít dòng thì tỉ lệ không đáng tin, bỏ qua kiểm tra


def _chuan_hoa_nfc(text: str) -> str:
    # Chuẩn hóa Unicode NFC ngay khi đọc: file PDF/PPTX có thể lưu tiếng Việt ở dạng NFD
    # (tổ hợp dấu rời), nếu không chuẩn hóa sẽ khiến 2 chuỗi "giống nhau" bị so khớp/tìm
    # kiếm sai lệch ở các bước sau (chunking, hiển thị trích dẫn).
    return unicodedata.normalize("NFC", text)


def _don_dep_watermark(text: str) -> str:
    """Dọn nhiễu lặp lại trong text vừa đọc: watermark và ký tự PDF không giải mã được."""
    text = _MAU_WATERMARK_STUDOCU.sub("", text)
    # Xoá mã ký tự PDF chưa giải được. Khi PDF nhúng font không kèm bảng ánh xạ ToUnicode
    # (rất hay gặp với font toán học), pdfplumber không biết ký tự đó là chữ gì nên trả về
    # nguyên mã dạng "(cid:10)". Đo trên giáo trình Bishop: 1761/4197 chunk (42%) dính loại
    # rác này. Chúng vô nghĩa với người đọc lẫn với embedding, nhưng vẫn chiếm token trong
    # chunk và làm loãng vector - xoá đi là dọn nhiễu, KHÔNG mất thông tin (thông tin đó đã
    # mất từ khâu đọc file, không có cách nào khôi phục nếu không dùng thêm OCR).
    text = _MAU_CID_PDF.sub("", text)
    # Xoá đi để lại khoảng trắng thừa (nhiều (cid:) liên tiếp) -> gộp lại cho gọn.
    return _MAU_KHOANG_TRANG_THUA.sub(" ", text)


def _la_trang_muc_luc(text: str) -> bool:
    """Đoán 1 trang có phải Mục lục (hoặc là 1 trang TIẾP THEO của mục lục nhiều trang) hay
    không. 2 cách nhận diện, chỉ cần khớp 1:

    1. Trang có tiêu đề "Mục lục"/"Table of Contents" ở ~40 ký tự đầu - bắt được trang ĐẦU
       TIÊN của mục lục. Chỉ xét đầu trang để tránh báo nhầm 1 trang nội dung chỉ tình cờ
       nhắc tới cụm từ này ở giữa bài.
    2. Đa số dòng trong trang kết thúc bằng 1 con số (dạng "<tên mục>... <số trang>") - bắt
       được các trang TIẾP THEO của mục lục nhiều trang, vốn không lặp lại tiêu đề nên cách
       (1) bỏ sót (đã gặp thực tế: mục lục chi tiết của giáo trình dài tới 6 trang, chỉ
       trang đầu có chữ "Mục lục").

    Cả 2 đều là suy đoán xấp xỉ có chủ đích (không phân tích cấu trúc PDF thật sự, ví dụ
    không dùng outline/bookmark của file) - ngưỡng của cách (2) đã hiệu chỉnh bằng dữ liệu
    thật, xem ARCHITECTURE.md.
    """
    dau_trang = text.strip().lower()[:40]
    if any(dau_trang.startswith(tieu_de) for tieu_de in _CAC_TIEU_DE_MUC_LUC):
        return True

    cac_dong = [dong for dong in text.split("\n") if dong.strip()]
    if len(cac_dong) < _SO_DONG_TOI_THIEU_DE_XET:
        return False
    so_dong_ket_thuc_bang_so = sum(1 for dong in cac_dong if _MAU_DONG_KET_THUC_BANG_SO.search(dong))
    return (so_dong_ket_thuc_bang_so / len(cac_dong)) >= _NGUONG_TY_LE_MUC_LUC


def _danh_dau_tieu_de(text: str, cap: int = 2) -> str:
    """Bọc 1 dòng thành dấu tiêu đề kiểu Markdown để splitter ưu tiên cắt tại đây.

    Dùng ký hiệu Markdown (`#`, `##`) thay vì mốc riêng như [BẢNG] vì 2 lẽ: LLM hiểu ngay
    đây là tiêu đề (định dạng cực phổ biến), và Streamlit render nó thành tiêu đề thật khi
    hiển thị trích dẫn - không phải viết thêm code cho cả hai việc.
    """
    return f"\n{'#' * cap} {text.strip()}\n"


def _phat_hien_tieu_de_pdf(trang) -> set:
    """Đoán những dòng nào trong trang PDF là tiêu đề, dựa vào cỡ chữ và độ dài.

    PDF không lưu cấu trúc logic: không có khái niệm "Heading 1" như DOCX, không có
    placeholder tiêu đề như PPTX. Thứ duy nhất còn lại là HÌNH THỨC - tiêu đề thường được
    in to hơn phần thân và ngắn hơn một câu văn.

    Lấy cỡ chữ áp đảo của trang (mode) làm mốc so sánh thay vì một con số cứng, vì mỗi tài
    liệu dùng cỡ chữ khác nhau (11pt, 12pt, 13pt...) - so theo TỈ LỆ mới tổng quát được.

    Đây là suy đoán xấp xỉ có chủ đích, giống cách nhận diện trang Mục lục ở §5.17: có thể
    bỏ sót hoặc nhận nhầm, nhưng sai ở đây chỉ làm chunk bị cắt hơi khác đi, không làm mất
    nội dung - nên đánh đổi chấp nhận được so với việc thêm một model phân tích bố cục.
    """
    cac_ky_tu = trang.chars
    if not cac_ky_tu:
        return set()

    # Gom ký tự thành dòng theo toạ độ "top" (làm tròn để chịu được sai lệch nhỏ khi render).
    theo_dong: Dict[int, List] = {}
    for ky_tu in cac_ky_tu:
        theo_dong.setdefault(round(ky_tu["top"]), []).append(ky_tu)

    dem_co_chu = Counter(round(k["size"], 1) for k in cac_ky_tu)
    co_chu_ap_dao = dem_co_chu.most_common(1)[0][0]
    if co_chu_ap_dao <= 0:
        return set()

    tieu_de = set()
    for cac_ky_tu_dong in theo_dong.values():
        noi_dung = "".join(k["text"] for k in cac_ky_tu_dong).strip()
        if not noi_dung or len(noi_dung) > config.DO_DAI_TOI_DA_TIEU_DE:
            continue
        co_chu_dong = sum(k["size"] for k in cac_ky_tu_dong) / len(cac_ky_tu_dong)
        if co_chu_dong >= co_chu_ap_dao * config.TY_LE_KICH_THUOC_CHU_TIEU_DE:
            tieu_de.add(noi_dung)
    return tieu_de


def _o_bang_khong_lap(hang) -> List[str]:
    """Lấy text các ô của 1 hàng, KHỬ phần nhân bản do ô gộp (merged cell).

    python-docx/python-pptx trả về CÙNG một ô cho mọi cột mà ô gộp trải qua. Với biểu mẫu
    hành chính - vốn gộp ô rất nhiều - điều này khiến cùng một chuỗi lặp lại hàng chục lần.
    Đo trên một file đăng ký đề tài NCKH thật (bảng 24x11): đọc thô cho ra 131.115 ký tự cho
    một biểu mẫu vài trang, tức nội dung bị nhân lên khoảng 10 lần. Hậu quả không chỉ là
    phình index mà là hỏng retrieval: chunk nào cũng na ná chunk nào nên không đoạn nào nổi
    bật lên được, và ngân sách ngữ cảnh bị đốt vào nội dung trùng.

    Nhận diện ô gộp bằng ĐỐI TƯỢNG XML nền (`_tc`/`_tbl`) chứ không so text: hai ô khác nhau
    hoàn toàn có thể tình cờ trùng nội dung ("Có"/"Có"), lúc đó xoá đi là mất dữ liệu thật.
    """
    ket_qua, da_thay = [], set()
    for o in hang.cells:
        # Mỗi thư viện đặt tên phần tử XML nền khác nhau; lấy cái nào có.
        # Dùng `is not None` chứ KHÔNG dùng `or`: phần tử lxml hưởng ứng truth-test theo số
        # phần tử con, nên một ô rỗng (không có con) sẽ bị coi là falsy và rơi nhầm sang
        # nhánh sau - lúc đó khoá so trùng lấy từ đối tượng khác, việc khử ô gộp hỏng âm
        # thầm. lxml cũng cảnh báo chính điều này (FutureWarning: truth-testing).
        nen = getattr(o, "_tc", None)
        if nen is None:
            nen = getattr(o, "_element", None)
        khoa = id(nen) if nen is not None else None
        if khoa is not None and khoa in da_thay:
            continue  # cùng một ô gộp trải qua nhiều cột -> chỉ lấy 1 lần
        if khoa is not None:
            da_thay.add(khoa)
        ket_qua.append(o.text)
    return ket_qua


def _bang_sang_markdown(bang: List[List]) -> str:
    """Đổi bảng (list hàng × ô) thành bảng Markdown.

    Dùng chung cho cả 3 định dạng vì pdfplumber/python-docx/python-pptx đều trả bảng về
    cùng một hình dạng list-of-rows. Chọn Markdown thay vì text phẳng vì 3 lý do:
      - Giữ được quan hệ HÀNG-CỘT: "5.000 đồng" thuộc dòng "Tạp chí đóng tập" chứ không
        phải một con số trôi nổi. Làm phẳng thành dòng rời là mất chính thông tin đó.
      - LLM đọc Markdown table rất tốt (định dạng phổ biến trong dữ liệu huấn luyện).
      - Streamlit render thẳng thành bảng đẹp ở phần trích dẫn, không cần code thêm.
    """
    cac_hang = [
        ["" if o is None else " ".join(str(o).split()) for o in hang]
        for hang in bang
        if hang and any(o is not None and str(o).strip() for o in hang)
    ]
    if not cac_hang:
        return ""
    so_cot = max(len(h) for h in cac_hang)
    cac_hang = [h + [""] * (so_cot - len(h)) for h in cac_hang]

    # Bỏ hẳn những cột RỖNG Ở MỌI HÀNG. Biểu mẫu Word hay được kẻ dư cột để căn lề, nên
    # bảng thật 3 cột lại ra Markdown 5 cột với 2 cột trắng - mỗi hàng gánh thêm "| |"
    # không mang thông tin nào mà vẫn tốn token của chunk và làm loãng vector. Cột rỗng ở
    # mọi hàng thì chắc chắn không mất mát gì khi bỏ.
    cot_co_chu = [i for i in range(so_cot) if any(h[i].strip() for h in cac_hang)]
    if cot_co_chu and len(cot_co_chu) < so_cot:
        cac_hang = [[h[i] for i in cot_co_chu] for h in cac_hang]
        so_cot = len(cot_co_chu)

    dong = ["| " + " | ".join(cac_hang[0]) + " |",
            "| " + " | ".join(["---"] * so_cot) + " |"]
    dong += ["| " + " | ".join(h) + " |" for h in cac_hang[1:]]
    return "\n".join(dong)


def _lay_client_vision():
    """Tạo Ollama client dùng chung cho OCR, chỉ một lần cho cả lần build.

    Trả về None (kèm cảnh báo một lần) nếu model vision chưa được pull - lúc đó OCR tự bỏ
    qua, phần còn lại của luồng đọc tài liệu vẫn chạy bình thường.
    """
    global _client_vision, _da_canh_bao_vision
    if _client_vision is None and not _da_canh_bao_vision:
        import ollama

        client = ollama.Client(host=config.OLLAMA_HOST)
        if mo_hinh_vision_co_san(client):
            _client_vision = client
        else:
            _da_canh_bao_vision = True
            logger.warning(
                "BAT_OCR_DU_PHONG đang bật nhưng model vision '%s' chưa được pull - bỏ qua "
                "OCR. Chạy: ollama pull %s",
                config.VISION_MODEL_NAME, config.VISION_MODEL_NAME,
            )
    return _client_vision


def _ocr_cac_trang(
    pdf, cac_so_trang: List[int], ten_file: str, bam_tai_lieu: str
) -> Dict[int, str]:
    """OCR một LOẠT trang PDF cùng lúc, trả về {số trang: text đọc được}.

    Ba thay đổi so với bản cũ (OCR từng trang ngay trong vòng lặp đọc), tất cả đều nhắm vào
    cùng một chỗ đau: đây là bước đắt nhất của luồng Ingestion với tài liệu scan.

      1. TRA CACHE TRƯỚC KHI RENDER. Khoá cache là (băm nội dung tài liệu, số trang, DPI) chứ
         không phải ảnh đã render - nên khi trúng cache thì KHÔNG phải render gì cả. Render
         một trang ở 150 DPI mất 0,2-0,4 giây; với sách scan 400 trang, khoá theo ảnh sẽ
         khiến cache chỉ tiết kiệm được một nửa chi phí (xem bo_nho_dem.khoa_ocr).
      2. GỌI SONG SONG. Mỗi lượt OCR là một yêu cầu HTTP tới Ollama rồi ngồi chờ; chạy
         nhiều lượt cùng lúc không tốn thêm CPU phía Python. Số worker suy từ phần cứng
         đang có, không phải một hằng số (xem tai_nguyen_gpu.so_worker_vision).
      3. RENDER VẪN TUẦN TỰ. pypdfium2 (bộ render nền của pdfplumber) không cam kết an toàn
         đa luồng, và render là việc thuần CPU nên thread cũng không giúp được gì. Chỉ phần
         CHỜ MẠNG mới được song song hoá - đúng chỗ có lợi, không hơn.

    Ảnh tạm bị xoá ngay sau khi dùng: với sách vài trăm trang, giữ lại toàn bộ ảnh render
    sẽ ngốn hàng GB đĩa mà không dùng vào việc gì (khác với ảnh TRÍCH TỪ tài liệu ở
    image_extractor - những ảnh đó còn phải hiển thị kèm trích dẫn).
    """
    if not cac_so_trang:
        return {}

    ket_qua: Dict[int, str] = {}
    can_goi: List[tuple] = []  # (so_trang, đường dẫn ảnh tạm, khoá cache)
    dung_cache = config.BAT_CACHE_INGESTION and bool(bam_tai_lieu)

    client = _lay_client_vision()
    for so_trang in cac_so_trang:
        khoa = bo_nho_dem.khoa_ocr(bam_tai_lieu, so_trang) if dung_cache else None
        if khoa:
            da_co = bo_nho_dem.kho_ocr.lay_text(khoa)
            if da_co is not None:
                ket_qua[so_trang] = da_co
                continue
        if client is None:
            continue  # model vision chưa sẵn sàng - đã cảnh báo một lần ở _lay_client_vision
        duong_dan_tam = config.IMAGES_DIR / f"_ocr_tam_{so_trang}.png"
        try:
            with do_thoi_gian.do("ocr_render_trang"):
                trang = pdf.pages[so_trang - 1]
                trang.to_image(resolution=config.DPI_RENDER_TRANG_OCR).original.save(duong_dan_tam)
        except Exception as loi:  # noqa: BLE001 - trang hỏng không được làm sập cả lần build
            logger.warning("Không render được trang %d của '%s': %s", so_trang, ten_file, loi)
            duong_dan_tam.unlink(missing_ok=True)
            continue
        can_goi.append((so_trang, duong_dan_tam, khoa))

    if not can_goi:
        if ket_qua:
            logger.info("OCR: cả %d trang đều lấy lại được từ cache.", len(ket_qua))
        return ket_qua

    so_worker = max(1, min(tai_nguyen_gpu.so_worker_vision(), len(can_goi)))
    logger.info(
        "OCR '%s': %d trang cần đọc lại bằng model vision (%d trang lấy từ cache), %d luồng.",
        ten_file, len(can_goi), len(ket_qua), so_worker,
    )

    def _chay(viec):
        so_trang, duong_dan_tam, khoa = viec
        try:
            text = ocr_trang_pdf(client, str(duong_dan_tam))
        except Exception as loi:  # noqa: BLE001
            logger.warning("Không OCR được trang %d của '%s': %s", so_trang, ten_file, loi)
            text = ""
        finally:
            duong_dan_tam.unlink(missing_ok=True)
        # Chỉ ghi cache khi ĐỌC ĐƯỢC. Ghi cả kết quả rỗng sẽ đóng băng một lần thất bại tạm
        # thời (Ollama đang tắt, model chưa nạp xong) thành kết quả vĩnh viễn.
        if text and khoa:
            bo_nho_dem.kho_ocr.luu_text(khoa, text)
        return so_trang, text

    with do_thoi_gian.do("ocr_goi_model"):
        if so_worker == 1:
            cac_ket_qua = [_chay(v) for v in can_goi]
        else:
            with ThreadPoolExecutor(max_workers=so_worker) as bo_chay:
                cac_ket_qua = list(bo_chay.map(_chay, can_goi))

    so_doc_duoc = 0
    for so_trang, text in cac_ket_qua:
        if text:
            ket_qua[so_trang] = text
            so_doc_duoc += 1
    logger.info(
        "Đã OCR lại %d/%d trang của '%s'.", so_doc_duoc, len(can_goi), ten_file,
    )
    return ket_qua


def _la_bang_that(bang: List[List]) -> bool:
    """Lọc "bảng" do thuật toán dò nhầm ra khỏi bảng thật.

    pdfplumber dò bảng bằng các đường kẻ trên trang. Trong slide bài giảng, khung viền
    trang trí và gạch chân tiêu đề cũng là đường kẻ - nên mỗi slide thường bị dò ra một
    "bảng" 2x1 chứa đúng dòng tiêu đề. Đo trên bộ slide Computer Vision thật: gần như mọi
    trang đều dính ("COMPUTER VISION", "CONTENTS", "LECTURE OUTCOMES"...).

    Đây không phải lỗi vô hại: _text_pdf_khong_ke_bang() loại vùng bảng ra khỏi luồng text
    thường, nên tiêu đề slide - phần mang thông tin đậm đặc nhất - bị bóc khỏi văn bản rồi
    bọc lại thành "bảng" một cột vô nghĩa, đồng thời mất luôn cơ hội được nhận diện là tiêu
    đề (§5.25).

    Điều kiện tối thiểu để coi là bảng thật: phải có ÍT NHẤT 2 CỘT và 2 HÀNG có nội dung.
    Bảng 1 cột không mang quan hệ hàng-cột nào - thứ duy nhất khiến một bảng đáng được giữ
    nguyên khối - nên để nó chảy vào text thường vẫn tốt hơn.
    """
    cac_hang = [h for h in bang if h and any(o is not None and str(o).strip() for o in h)]
    if len(cac_hang) < 2:
        return False
    so_cot_co_chu = max(
        sum(1 for o in h if o is not None and str(o).strip()) for h in cac_hang
    )
    return so_cot_co_chu >= 2


def _khoi_bang(bang: List[List]) -> str:
    """Bọc 1 bảng trong cặp mốc để chunking nhận ra và giữ nguyên khối.

    Trả về chuỗi rỗng với bảng dò nhầm (xem _la_bang_that) - chỗ gọi sẽ để nội dung đó
    chảy vào luồng text thường thay vì bọc thành bảng.
    """
    if not _la_bang_that(bang):
        return ""
    markdown = _bang_sang_markdown(bang)
    return f"\n\n{MOC_BANG_MO}\n{markdown}\n{MOC_BANG_DONG}\n" if markdown else ""


def _text_pdf_khong_ke_bang(
    trang, cac_bang, ten_file: str = "", hieu_chinh: Optional[HieuChinhXTolerance] = None
) -> str:
    """Lấy văn xuôi của trang, ĐÃ LOẠI vùng chiếm bởi bảng.

    Nếu không loại, cùng một bảng sẽ vào index hai lần dưới hai dạng: bản Markdown sạch và
    bản text thô do extract_text() đọc ngang qua các ô (các cột dính vào nhau thành câu vô
    nghĩa). Bản thô đó vừa gây nhiễu retrieval vừa có thể lọt vào ngữ cảnh gửi cho LLM.

    outside_bbox() có thể ném lỗi khi bbox của bảng chạm mép trang; khi đó lùi về đọc cả
    trang - thà chấp nhận trùng lặp còn hơn mất trắng nội dung trang.
    """
    so_trang = getattr(trang, "page_number", "?")
    try:
        vung = trang
        for bang in cac_bang:
            vung = vung.outside_bbox(bang.bbox)
        return _trich_text_thich_ung(vung, ten_file, so_trang, hieu_chinh)
    except (ValueError, TypeError) as loi:
        logger.warning(
            "Không loại được vùng bảng ở trang %s (%s) - đọc cả trang, bảng có thể bị lặp.",
            so_trang, type(loi).__name__,
        )
        return _trich_text_thich_ung(trang, ten_file, so_trang, hieu_chinh)


def _cac_cot_cua_trang(trang) -> List[tuple]:
    """Trả về danh sách khoảng (x_trái, x_phải) của từng cột; rỗng nếu trang MỘT cột.

    VẤN ĐỀ: pdfplumber đọc theo DÒNG NGANG chạy suốt bề ngang trang. Với trang chia 2 cột,
    nó ghép câu của cột trái với câu của cột phải thành một dòng. Đo trên PDF 2 cột dựng thử:
        "Dieu 1. Pham vi dieu chinh cua luat nay Dieu 2. Doi tuong ap dung bao gom moi"
    Hai điều luật khác nhau dính thành một câu. Mọi chunk sinh ra từ trang đó đều vô nghĩa.

    CÁCH DÒ: chiếu tất cả các từ lên trục ngang, chia thành SO_O_DO_COT dải, rồi tìm dải
    trống liên tiếp - đó là rãnh giữa hai cột. Đây là dấu hiệu hình học nên không phụ thuộc
    ngôn ngữ hay loại tài liệu.

    HAI CHỐT CHỐNG BÁO ĐỘNG GIẢ, cả hai đều đến từ đo đạc chứ không phải phòng xa. Bản dò
    đầu tiên thiếu chúng và nhận nhầm 100% số trang của giáo trình Pháp luật thành 2 cột -
    tức sẽ phá nát một tài liệu đang đọc tốt:

      1. Rãnh phải nằm TRONG KHỐI CHỮ. Thứ bản đầu tưởng là rãnh thực ra là LỀ TRANG (dải
         trống ở rìa trái/phải). Nay chỉ xét khoảng từ dải có chữ đầu tiên tới dải có chữ
         cuối cùng.
      2. Hai bên rãnh đều phải có LƯỢNG CHỮ ĐÁNG KỂ (>= TY_LE_TU_MOI_COT). Một trang một cột
         có hình chèn giữa cũng tạo ra khoảng trống dọc, nhưng khi đó một bên gần như không
         có chữ.

    Đo lại sau khi có hai chốt này: 0 báo nhầm trên toàn bộ 9 file PDF của corpus, trong khi
    vẫn bắt đúng trang 2 cột dựng thử.
    """
    try:
        cac_tu = trang.extract_words()
    except Exception:  # noqa: BLE001 - trang hỏng không được làm sập cả lần đọc tài liệu
        return []
    rong = float(trang.width or 0)
    if len(cac_tu) < config.SO_TU_TOI_THIEU_DE_DO_COT or rong <= 0:
        return []

    so_o = config.SO_O_DO_COT
    co_chu = [False] * so_o
    for t in cac_tu:
        dau = max(0, min(so_o - 1, int(t["x0"] / rong * so_o)))
        cuoi = max(0, min(so_o - 1, int(t["x1"] / rong * so_o)))
        for i in range(dau, cuoi + 1):
            co_chu[i] = True

    o_co_chu = [i for i, x in enumerate(co_chu) if x]
    if not o_co_chu or o_co_chu[-1] - o_co_chu[0] < 12:
        return []

    cac_ranh, dau = [], None
    for i in range(o_co_chu[0] + 1, o_co_chu[-1]):
        if not co_chu[i]:
            dau = i if dau is None else dau
        else:
            if dau is not None and i - dau >= config.SO_O_RANH_TOI_THIEU:
                cac_ranh.append((dau + i) / 2 * rong / so_o)
            dau = None

    moc = []
    for giua in cac_ranh:
        ben_trai = sum(1 for t in cac_tu if t["x1"] <= giua)
        ben_phai = sum(1 for t in cac_tu if t["x0"] >= giua)
        if min(ben_trai, ben_phai) >= config.TY_LE_TU_MOI_COT * len(cac_tu):
            moc.append(giua)
    if not moc:
        return []

    ranh_gioi = [0.0] + moc + [rong]
    return [(ranh_gioi[i], ranh_gioi[i + 1]) for i in range(len(ranh_gioi) - 1)]


def _text_theo_cot(
    trang, cac_cot: List[tuple], ten_file: str, so_trang: int,
    hieu_chinh: Optional[HieuChinhXTolerance] = None,
) -> str:
    """Đọc từng cột riêng rồi nối lại theo thứ tự trái -> phải (thứ tự đọc của người)."""
    cac_phan = []
    for x0, x1 in cac_cot:
        try:
            vung = trang.crop((x0, trang.bbox[1], x1, trang.bbox[3]))
        except (ValueError, TypeError):
            continue
        phan = _trich_text_thich_ung(vung, ten_file, so_trang, hieu_chinh).strip()
        if phan:
            cac_phan.append(phan)
    return "\n\n".join(cac_phan)


def _doc_mot_trang_pdf(trang, ten_file: str, so_trang: int, hieu_chinh) -> str:
    """Toàn bộ việc đọc TEXT của một trang PDF, gói lại thành đúng MỘT lượt.

    Ba đường đọc (có bảng / nhiều cột / thường) loại trừ nhau, nên một trang chỉ đi qua đúng
    một đường. Thay đổi so với bản trước tuy nhỏ mà đáng kể: kết quả `extract()` của mỗi
    bảng được dùng LẠI thay vì gọi lần thứ hai lúc dựng khối bảng. Trích xuất một bảng là
    thao tác đắt (dò đường kẻ, gom ô, ghép text từng ô), và nó đang chạy hai lần cho mọi
    bảng của mọi trang chỉ vì hai chỗ gọi nằm cách nhau ba dòng.
    """
    # Lọc bảng dò nhầm NGAY TẠI ĐÂY, trước khi loại vùng bảng khỏi text thường. Nếu lọc muộn
    # hơn (lúc dựng khối bảng), phần text nằm trong vùng bảng giả đã bị
    # _text_pdf_khong_ke_bang() bóc đi rồi mà không được bọc lại - tức MẤT HẲN nội dung, tệ
    # hơn hẳn so với việc chỉ bọc sai định dạng.
    cac_bang = []
    for bang in trang.find_tables():
        cac_hang = bang.extract()
        if _la_bang_that(cac_hang):
            cac_bang.append((bang, cac_hang))

    if cac_bang:
        noidung = _text_pdf_khong_ke_bang(
            trang, [b for b, _ in cac_bang], ten_file, hieu_chinh
        )
        return noidung + "".join(_khoi_bang(cac_hang) for _, cac_hang in cac_bang)

    # Dò bố cục nhiều cột TRƯỚC khi đọc: đọc ngang một trang 2 cột sẽ trộn câu của hai cột
    # vào nhau, hỏng ngay từ khâu đầu tiên (xem _cac_cot_cua_trang). Chỉ làm khi trang không
    # có bảng - bảng vốn đã được tách vùng riêng, và một bảng nhiều cột sẽ khiến phép dò rãnh
    # hiểu nhầm thành bố cục nhiều cột.
    cac_cot = _cac_cot_cua_trang(trang) if config.BAT_DOC_THEO_COT else []
    if cac_cot:
        logger.info(
            "Trang %d của '%s' có bố cục %d cột - đọc từng cột riêng.",
            so_trang, ten_file, len(cac_cot),
        )
        return _text_theo_cot(trang, cac_cot, ten_file, so_trang, hieu_chinh)
    return _trich_text_thich_ung(trang, ten_file, so_trang, hieu_chinh)


def doc_pdf(duong_dan: Path) -> List[Dict]:
    """Đọc từng trang PDF bằng pdfplumber, trả về list dict {nguon, trang, noidung}.

    LUỒNG MỘT-LƯỢT-DUYỆT, chia 4 pha. Bản trước duyệt toàn bộ PDF HAI lần (một lần đọc text,
    một lần nữa trong trich_anh_pdf để trích ảnh - lần thứ hai còn gọi lại extract_text() cho
    từng trang chỉ để lấy dòng chú thích), và trong lần thứ nhất thì mỗi trang có thể bị đọc
    tới 5 lần vì phép dò x_tolerance. Với giáo trình Bishop, đó là gốc rễ của việc "cùng một
    tài liệu bị quét nhiều lần" mà log đã cho thấy.

      Pha 1 (tuần tự, MỘT lượt qua các trang): đọc text, dò bảng/cột, nhận diện tiêu đề,
             liệt kê ứng viên ảnh (chưa render), đánh dấu trang cần OCR. Sau mỗi trang gọi
             flush_cache() để nhả bộ nhớ - một cuốn sách 700 trang giữ hết đối tượng của mọi
             trang trong RAM là chuyện không cần thiết.
      Pha 2: OCR các trang đã đánh dấu - tra cache trước, phần còn lại gọi model song song.
      Pha 3: gộp kết quả OCR, dọn watermark, chuẩn hoá, loại trang rỗng / trang mục lục.
      Pha 4: render và lưu những ảnh được giữ lại (chỉ những trang thật sự có ảnh mới phải
             đọc lại), rồi loại các hình lặp lại kiểu logo/watermark.

    Thứ tự này giữ NGUYÊN mọi quy tắc nội dung của bản trước - kể cả quy tắc tinh tế nhất:
    ảnh phủ kín một trang chỉ bị loại khi OCR đã CHỨNG MINH trang đó là ảnh chụp một trang
    chữ, chứ không phải khi đoán qua kích thước (xem _la_anh_cua_trang_chu).
    """
    ten_file = duong_dan.name
    hieu_chinh = HieuChinhXTolerance()
    bam_tai_lieu = (
        bo_nho_dem.bam_file(duong_dan) if config.BAT_CACHE_INGESTION else ""
    )
    cac_trang_tho: List[Dict] = []
    ung_vien_anh: Dict[int, list] = {}

    with pdfplumber.open(duong_dan) as pdf:
        # ---------- PHA 1: đọc text, một lượt duy nhất ----------
        for so_trang, trang in enumerate(pdf.pages, start=1):
            with do_thoi_gian.do("pdf_doc_text_trang"):
                noidung = _doc_mot_trang_pdf(trang, ten_file, so_trang, hieu_chinh)

            # OCR DỰ PHÒNG: chỉ chạy khi việc trích xuất text đã THẤT BẠI (chữ ra thành mã
            # (cid:NN) do font thiếu bảng ánh xạ, hoặc trang scan không có chữ). Quyết định ở
            # đây - sau khi đã có text để đánh giá, trước khi dọn dẹp - vì phải nhìn thấy mã
            # (cid:) nguyên vẹn mới nhận ra được là trang hỏng.
            can_ocr = config.BAT_OCR_DU_PHONG and trang_can_ocr(noidung, len(trang.images))

            if config.BAT_NHAN_DIEN_TIEU_DE:
                with do_thoi_gian.do("pdf_nhan_dien_tieu_de"):
                    cac_tieu_de = _phat_hien_tieu_de_pdf(trang)
                for dong_tieu_de in cac_tieu_de:
                    # Thay đúng dòng đó bằng bản có dấu tiêu đề. Dùng thay-thế-theo-dòng
                    # thay vì dựng lại cả trang để không phá thứ tự nội dung gốc.
                    noidung = noidung.replace(dong_tieu_de, _danh_dau_tieu_de(dong_tieu_de), 1)

            if config.BAT_TRICH_ANH:
                cac_ung_vien = ung_vien_anh_trang(trang)
                if cac_ung_vien:
                    ung_vien_anh[so_trang] = cac_ung_vien

            cac_trang_tho.append(
                {"so_trang": so_trang, "noidung": noidung, "can_ocr": can_ocr}
            )
            # Nhả các đối tượng đã phân tích của trang này. Pha 4 chỉ đọc lại đúng những
            # trang có ảnh cần render, nên với tài liệu thuần chữ thì không trang nào phải
            # phân tích lần thứ hai.
            trang.flush_cache()

        # ---------- PHA 2: OCR song song những trang đã đánh dấu ----------
        ket_qua_ocr = _ocr_cac_trang(
            pdf, [t["so_trang"] for t in cac_trang_tho if t["can_ocr"]], ten_file, bam_tai_lieu
        )

        # ---------- PHA 3: gộp OCR + dọn dẹp + lọc ----------
        ket_qua: List[Dict] = []
        # Những trang mà OCR đã lấy ra được CHỮ THẬT. Dùng để quyết định có trích ảnh của
        # trang đó nữa hay không: nếu OCR đọc ra cả một trang chữ thì "ảnh" của trang chính
        # là ảnh chụp trang chữ ấy - trích thêm nó ra chỉ tạo một chunk hình vô nghĩa cho mỗi
        # trang. Đây là tín hiệu ĐO ĐƯỢC, thay cho việc đoán qua kích thước ảnh.
        cac_trang_ocr_ra_chu = set()
        noi_dung_theo_trang: Dict[int, str] = {}
        for tho in cac_trang_tho:
            so_trang = tho["so_trang"]
            noidung = tho["noidung"]
            noidung_ocr = ket_qua_ocr.get(so_trang, "")
            if noidung_ocr:
                noidung = noidung_ocr
                if len(noidung_ocr.split()) >= config.SO_TU_TOI_THIEU_TRANG_CO_CHU:
                    cac_trang_ocr_ra_chu.add(so_trang)

            noidung = _don_dep_watermark(noidung)
            noidung = _chuan_hoa_nfc(noidung).strip()
            noi_dung_theo_trang[so_trang] = noidung
            if not noidung:
                logger.warning(
                    "Trang %d của '%s' không có text (có thể là ảnh scan) - bỏ qua.",
                    so_trang, ten_file,
                )
                continue
            if _la_trang_muc_luc(noidung):
                logger.warning(
                    "Trang %d của '%s' có vẻ là trang Mục lục - bỏ qua (không phải nội "
                    "dung thật, dễ gây nhiễu retrieval).",
                    so_trang, ten_file,
                )
                continue
            ket_qua.append({"nguon": ten_file, "trang": so_trang, "noidung": noidung})

        # ---------- PHA 4: render những ảnh được giữ lại ----------
        if config.BAT_TRICH_ANH and ung_vien_anh:
            with do_thoi_gian.do("pdf_trich_anh"):
                cac_ban_ghi_anh, so_anh_toan_trang = [], 0
                for so_trang, cac_ung_vien in ung_vien_anh.items():
                    la_trang_scan_chu = so_trang in cac_trang_ocr_ra_chu
                    cac_bbox = []
                    for bbox, phu_ca_trang in cac_ung_vien:
                        if la_trang_scan_chu and phu_ca_trang:
                            # Ảnh chụp CẢ TRANG của một PDF scan: nội dung thật của nó là
                            # CHỮ, và chữ đó đã được OCR lấy ra ở pha 2. Trích thêm một
                            # "hình" nữa chỉ tạo một chunk rác cho mỗi trang, tốn thêm một
                            # lượt model vision cho mỗi trang, và khiến trích dẫn trỏ vào
                            # "hình" thay vì vào đoạn chữ mà câu trả lời thật sự dùng.
                            so_anh_toan_trang += 1
                            continue
                        cac_bbox.append(bbox)
                    if not cac_bbox:
                        continue
                    cac_ban_ghi_anh.extend(
                        luu_anh_trang_pdf(
                            ten_file, pdf.pages[so_trang - 1], so_trang, cac_bbox,
                            noi_dung_theo_trang.get(so_trang, "").split("\n"),
                        )
                    )
                    pdf.pages[so_trang - 1].flush_cache()
                if so_anh_toan_trang:
                    logger.info(
                        "Bỏ qua %d ảnh chụp cả trang ở '%s' (PDF scan - nội dung của chúng "
                        "là chữ, đã được OCR đọc ra).", so_anh_toan_trang, ten_file,
                    )
                ket_qua.extend(loc_anh_lap_lai(cac_ban_ghi_anh, ten_file))
    return ket_qua


def duyet_shape(cac_shape) -> Iterator:
    """Duyệt shape của slide, ĐỆ QUY vào các nhóm (group shape).

    PowerPoint cho phép gộp nhiều hình thành 1 nhóm; khi đó `for shape in slide.shapes` chỉ
    thấy đúng cái nhóm, không thấy gì bên trong. Mọi ảnh/bảng/chữ nằm trong nhóm sẽ biến mất
    khỏi index mà KHÔNG có lỗi nào báo ra - đã đo trên tài liệu mẫu: một caption và một ảnh
    trong nhóm bị mất hoàn toàn. Hàm này là chỗ duy nhất xử lý việc đó, để cả phần đọc text
    lẫn phần trích ảnh dùng chung, không nơi nào quên đệ quy.
    """
    for shape in cac_shape:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from duyet_shape(shape.shapes)
        else:
            yield shape


def _trich_text_shape(shape) -> List[str]:
    """Lấy text từ 1 shape trong slide: text_frame và table (nếu có)."""
    cac_doan = []
    if shape.has_text_frame:
        for doan_van in shape.text_frame.paragraphs:
            text = "".join(run.text for run in doan_van.runs)
            if text.strip():
                cac_doan.append(text)
    if shape.has_table:
        # Giữ nguyên cấu trúc hàng-cột dưới dạng Markdown. Bản trước duyệt từng ô rồi nối
        # thành các dòng rời rạc, làm mất hoàn toàn quan hệ "giá trị này thuộc dòng nào,
        # cột nào" - câu hỏi kiểu "phí phạt của tạp chí đóng tập là bao nhiêu" vì thế không
        # thể trả lời đúng dù dữ liệu vẫn nằm trong index.
        cac_doan.append(
            _khoi_bang([_o_bang_khong_lap(hang) for hang in shape.table.rows])
        )
    return cac_doan


def doc_pptx(duong_dan: Path) -> List[Dict]:
    """Đọc từng slide PPTX bằng python-pptx, trả về list dict {nguon, trang, noidung}.

    "trang" ở đây là số thứ tự slide (để đồng nhất field với PDF, giúp citation.py
    dùng chung 1 tên field cho cả 2 loại tài liệu).
    """
    ket_qua = []
    trinh_chieu = Presentation(duong_dan)
    for so_slide, slide in enumerate(trinh_chieu.slides, start=1):
        # Tiêu đề slide là tín hiệu cấu trúc có sẵn trong định dạng (placeholder title),
        # chắc chắn hơn hẳn việc đoán từ cỡ chữ như với PDF.
        shape_tieu_de = slide.shapes.title if slide.shapes.title is not None else None
        cac_doan = []
        # So sánh bằng shape_id chứ KHÔNG bằng `is`: python-pptx dựng một object wrapper MỚI
        # mỗi lần truy cập cùng một shape, nên `slide.shapes.title is shape` luôn False dù
        # đúng là một shape - hệ quả là tiêu đề bị ghi 2 lần (1 lần dạng tiêu đề, 1 lần dạng
        # text thường), làm nhiễu chính đoạn quan trọng nhất của slide.
        id_tieu_de = None
        if config.BAT_NHAN_DIEN_TIEU_DE and shape_tieu_de is not None:
            van_ban_tieu_de = shape_tieu_de.text_frame.text.strip()
            if van_ban_tieu_de:
                cac_doan.append(_danh_dau_tieu_de(van_ban_tieu_de))
                id_tieu_de = shape_tieu_de.shape_id
        for shape in duyet_shape(slide.shapes):
            if id_tieu_de is not None and shape.shape_id == id_tieu_de:
                continue  # đã thêm ở trên dưới dạng tiêu đề, không lặp lại thành text thường
            cac_doan.extend(_trich_text_shape(shape))
        noidung = _don_dep_watermark("\n".join(cac_doan))
        noidung = _chuan_hoa_nfc(noidung).strip()
        if not noidung:
            logger.warning(
                "Slide %d của '%s' không có text (có thể chỉ chứa ảnh) - bỏ qua.",
                so_slide,
                duong_dan.name,
            )
            continue
        if _la_trang_muc_luc(noidung):
            logger.warning(
                "Slide %d của '%s' có vẻ là slide Mục lục - bỏ qua (không phải nội dung "
                "thật, dễ gây nhiễu retrieval).",
                so_slide,
                duong_dan.name,
            )
            continue
        ket_qua.append({"nguon": duong_dan.name, "trang": so_slide, "noidung": noidung})
    if config.BAT_TRICH_ANH:
        with do_thoi_gian.do("pptx_trich_anh"):
            ket_qua.extend(
                loc_anh_lap_lai(trich_anh_pptx(duong_dan, trinh_chieu), duong_dan.name)
            )
    return ket_qua


_NS_WORD = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
}


def _text_trong_text_box(doan_van) -> str:
    """Lấy chữ nằm trong TEXT BOX của một đoạn văn Word (python-docx không thấy được).

    `Paragraph.text` chỉ gom text của các run trực tiếp trong đoạn. Nội dung text box nằm
    sâu trong <w:drawing> -> <wps:txbx> -> <w:txbxContent> -> các <w:p> con, nên nó vô hình
    với python-docx và biến mất khỏi index mà không có bất kỳ dấu hiệu nào.

    Vì sao đáng sửa chứ không phải trường hợp hiếm: trong Word, sơ đồ, khung "Lưu ý", trích
    dẫn nổi bật, chú thích bên lề đều hay được đặt trong text box - và đó thường là chỗ chứa
    thông tin cô đọng nhất của trang. Mất chúng nghĩa là mất đúng phần đáng hỏi nhất.

    Dò bằng đường dẫn XML `.//w:txbxContent//w:t` thay vì theo namespace của riêng
    <wps:txbx>: text box do các phiên bản Word khác nhau tạo ra nằm dưới nhiều namespace
    (wps, v:textbox của VML cũ), nhưng tất cả đều chứa <w:txbxContent> bên trong.
    """
    try:
        cac_o = doan_van._element.findall(".//w:txbxContent//w:t", _NS_WORD)
    except Exception:  # noqa: BLE001 - XML lạ không được làm sập cả lần đọc tài liệu
        return ""
    cac_chu = [o.text for o in cac_o if o.text and o.text.strip()]
    return (" ".join(cac_chu) + "\n") if cac_chu else ""


def doc_docx(duong_dan: Path) -> List[Dict]:
    """Đọc file DOCX bằng python-docx, tách thành các "trang" theo dấu ngắt trang cứng
    (Insert > Page Break trong Word).

    Khác với PDF/PPTX, file .docx KHÔNG lưu số trang cố định trong XML - số trang thật khi
    in/xem phụ thuộc khổ giấy, lề, font lúc render, không phải thứ có sẵn trong file. Ngắt
    trang cứng là tín hiệu gần nhất có sẵn để tách nội dung thành các đơn vị tương tự
    "trang" mà không cần render tài liệu ra ảnh (ngoài phạm vi đồ án). Nếu tài liệu không
    có ngắt trang cứng nào, toàn bộ nội dung được coi là 1 "trang" duy nhất.

    Đọc cả đoạn văn LẪN bảng, theo ĐÚNG THỨ TỰ xuất hiện trong tài liệu, nhờ
    `document.iter_inner_content()` (trả lần lượt Paragraph | Table). Bản trước chỉ duyệt
    `document.paragraphs` nên bảng biến mất hoàn toàn - đã đo trên tài liệu mẫu: một file
    chứa 2 bảng chỉ đọc ra được 333 ký tự văn xuôi, toàn bộ số liệu trong bảng không hề vào
    index. Duyệt theo thứ tự tài liệu quan trọng hơn việc gom hết bảng về cuối, vì bảng phải
    nằm cạnh đoạn văn giới thiệu nó thì ngữ cảnh mới đúng.

    Giới hạn còn lại: bảng LỒNG trong ô của bảng khác vẫn không thấy (python-docx ghi rõ
    `document.tables` chỉ liệt kê bảng cấp cao nhất, và iter_inner_content cũng theo cấp
    đó). Chấp nhận được vì hiếm gặp; ghi ra đây để không ai tưởng đã xử lý xong.
    """
    document = Document(duong_dan)
    cac_trang_tho: List[str] = [""]
    for phan in document.iter_inner_content():
        if isinstance(phan, BangDocx):
            cac_trang_tho[-1] += _khoi_bang([_o_bang_khong_lap(hang) for hang in phan.rows])
            continue
        if not isinstance(phan, DoanVanDocx):
            continue
        if phan.text.strip():
            # style.name dạng "Heading 1"/"Heading 2"... là tín hiệu cấu trúc chuẩn của Word.
            ten_style = (phan.style.name or "") if phan.style is not None else ""
            if config.BAT_NHAN_DIEN_TIEU_DE and ten_style.startswith("Heading"):
                cap = int(ten_style.split()[-1]) if ten_style.split()[-1].isdigit() else 1
                cac_trang_tho[-1] += _danh_dau_tieu_de(phan.text, cap=min(cap, 2))
            else:
                cac_trang_tho[-1] += phan.text + "\n"
        # TEXT BOX: `phan.text` của python-docx KHÔNG bao gồm nội dung trong text box - nó
        # chỉ gom text của các run trực tiếp. Nội dung text box nằm sâu trong
        # <w:drawing>/<wps:txbx>/<w:txbxContent>, nên với python-docx nó vô hình và MẤT HẲN
        # khỏi index mà không có dấu hiệu gì. Đây là lỗ hổng thật với tài liệu chưa từng gặp:
        # sơ đồ, khung "Lưu ý", trích dẫn nổi bật trong Word đều hay được đặt trong text box,
        # và đó thường là chỗ chứa thông tin cô đọng nhất của trang.
        cac_trang_tho[-1] += _text_trong_text_box(phan)
        if any(run._element.xpath(".//w:br[@w:type='page']") for run in phan.runs):
            cac_trang_tho.append("")

    ket_qua = []
    for so_trang, noidung_tho in enumerate(cac_trang_tho, start=1):
        noidung = _don_dep_watermark(noidung_tho)
        noidung = _chuan_hoa_nfc(noidung).strip()
        if not noidung:
            logger.warning(
                "Trang %d của '%s' không có text - bỏ qua.", so_trang, duong_dan.name
            )
            continue
        if _la_trang_muc_luc(noidung):
            logger.warning(
                "Trang %d của '%s' có vẻ là trang Mục lục - bỏ qua (không phải nội dung "
                "thật, dễ gây nhiễu retrieval).",
                so_trang,
                duong_dan.name,
            )
            continue
        ket_qua.append({"nguon": duong_dan.name, "trang": so_trang, "noidung": noidung})
    if config.BAT_TRICH_ANH:
        with do_thoi_gian.do("docx_trich_anh"):
            ket_qua.extend(
                loc_anh_lap_lai(trich_anh_docx(duong_dan, document), duong_dan.name)
            )
    return ket_qua


def doc_tai_lieu(duong_dan: Path) -> List[Dict]:
    """Đọc 1 file, tự chọn hàm đọc phù hợp theo phần đuôi file."""
    duoi = duong_dan.suffix.lower()
    if duoi == ".pdf":
        return doc_pdf(duong_dan)
    if duoi == ".pptx":
        return doc_pptx(duong_dan)
    if duoi == ".docx":
        return doc_docx(duong_dan)
    raise ValueError(f"Định dạng không hỗ trợ: '{duoi}' (chỉ hỗ trợ {CAC_DUOI_HO_TRO})")


def doc_tai_lieu_hoan_chinh(duong_dan: Path) -> List[Dict]:
    """Đọc 1 file và làm TRỌN VẸN mọi bước sinh nội dung: đọc + chú thích ảnh + bỏ ảnh rỗng.

    Đây là ĐƠN VỊ ĐƯỢC CACHE, và ranh giới đó phải nằm đúng ở đây chứ không sớm hơn. Nếu chỉ
    cache kết quả đọc thô rồi vẫn chạy lại bước chú thích vision mỗi lần, ta sẽ cache đúng
    phần rẻ và bỏ qua phần đắt - benchmark của project cho thấy chú thích ảnh tốn ~1,9 giây
    mỗi hình, tức phần lớn thời gian của một lần build với tài liệu nhiều hình.

    Bước chú thích vision chuyển từ "gom toàn bộ ảnh của cả corpus rồi xử lý một lượt" sang
    "làm xong từng tài liệu": đó là điều kiện để kết quả của một tài liệu tự nó đầy đủ và
    cache được độc lập. Đổi lại, tiến độ nay báo theo từng tài liệu thay vì theo cả corpus -
    và đó lại là cách báo đúng hơn, vì với index tăng dần thì đa số tài liệu không được xử
    lý lại chút nào.
    """
    ket_qua = doc_tai_lieu(duong_dan)
    if config.BAT_TRICH_ANH and config.BAT_CHU_THICH_ANH:
        cac_anh = [m for m in ket_qua if m.get("loai_noi_dung") == "anh"]
        if cac_anh:
            logger.info("'%s': %d ảnh cần chú thích.", duong_dan.name, len(cac_anh))
            bo_sung_chu_thich_vision(cac_anh)
    return _bo_ban_ghi_anh_rong(ket_qua)


def _cache_con_du_anh(cac_trang: List[Dict]) -> bool:
    """Mọi file ảnh mà bản ghi trong cache trỏ tới còn tồn tại không?

    Cần kiểm tra vì hai phần dữ liệu nằm ở hai chỗ khác nhau: nội dung nằm trong cache
    (data/cache), còn file ảnh nằm ở data/images. Xoá data/images mà giữ cache sẽ cho ra một
    index "hợp lệ" với các trích dẫn trỏ vào ảnh không còn tồn tại - hỏng đúng ở chỗ người
    dùng nhìn thấy. Coi đó là cache trượt và đọc lại là cách xử lý rẻ và đúng.
    """
    return all(
        Path(m["duong_dan_anh"]).exists()
        for m in cac_trang
        if m.get("loai_noi_dung") == "anh" and m.get("duong_dan_anh")
    )


def doc_tai_lieu_co_cache(duong_dan: Path) -> List[Dict]:
    """doc_tai_lieu_hoan_chinh() nhưng lấy lại kết quả cũ khi nội dung file không đổi.

    Khoá cache = băm nội dung file + vân tay cấu hình đọc (xem rag/bo_nho_dem.py). Đây là
    thứ biến câu "tài liệu đã xử lý rồi thì không được xử lý lại" thành hành vi thật, kể cả
    khi người dùng đổi tên file, chuyển thư mục, hay build lại index vì một tài liệu KHÁC
    vừa thay đổi.
    """
    if not config.BAT_CACHE_INGESTION:
        return doc_tai_lieu_hoan_chinh(duong_dan)

    khoa = bo_nho_dem.khoa_tai_lieu(duong_dan)
    da_co = bo_nho_dem.kho_tai_lieu.lay_json(khoa)
    if da_co is not None and _cache_con_du_anh(da_co):
        logger.info(
            "'%s': lấy lại %d bản ghi từ cache (nội dung file không đổi).",
            duong_dan.name, len(da_co),
        )
        do_thoi_gian.ghi_nhan("tai_lieu_trung_cache", 0.0)
        return da_co

    with do_thoi_gian.do("tai_lieu_doc_moi"):
        ket_qua = doc_tai_lieu_hoan_chinh(duong_dan)
    bo_nho_dem.kho_tai_lieu.luu_json(khoa, ket_qua)
    return ket_qua


def _bo_ban_ghi_anh_rong(cac_trang: List[Dict]) -> List[Dict]:
    """Bỏ những bản ghi ảnh rốt cuộc KHÔNG có chú thích nào.

    Một ảnh không có chú thích lân cận và cũng không được model vision đọc sẽ cho ra nội
    dung đúng bằng chuỗi "[HÌNH]" - 6 ký tự, giống hệt nhau ở mọi ảnh, không mang một thông
    tin nào để tra cứu. Nhưng nó vẫn chiếm một vector trong index và vẫn lọt được vào top-K,
    tức vừa vô dụng vừa gây nhiễu.

    Đây là một lỗi thật, đo được: người dùng nạp vào một giáo trình 383 trang dạng SCAN (mỗi
    trang là một ảnh, không có lớp text nào) trong lúc Ollama đang tắt. Kết quả là index gồm
    ĐÚNG 379 chunk "[HÌNH]" giống hệt nhau; hệ thống báo "đã build index với 379 chunk" như
    thể thành công rồi trả lời sai mọi câu hỏi, vì nó thật sự không có gì để đọc. Không một
    exception nào được ném ra.

    Chunk văn xuôi vốn đã có ngưỡng độ dài tối thiểu chặn đúng loại rác này
    (chunking.DO_DAI_CHUNK_TOI_THIEU = 25 ký tự), nhưng bản ghi ảnh đi thẳng thành chunk nên
    lọt qua - đây là chỗ bịt lại, và phải nằm SAU bước chú thích vision mới đúng.
    """
    giu, bo = [], 0
    for m in cac_trang:
        if m.get("loai_noi_dung") == "anh" and m["noidung"].replace(MOC_ANH, "").strip() == "":
            bo += 1
            continue
        giu.append(m)
    if bo:
        logger.warning(
            "Bỏ %d ảnh không có chú thích nào (không có chữ lân cận, và model vision không "
            "chạy hoặc đang tắt) - chúng chỉ chứa mốc '%s' nên vô dụng cho tra cứu.",
            bo, MOC_ANH,
        )
    return giu


def _canh_bao_tai_lieu_khong_doc_duoc(cac_trang: List[Dict], cac_duong_dan: List[Path]) -> None:
    """Cảnh báo to khi một tài liệu gần như không đọc được chữ nào.

    Vì sao cần: build index "thành công" với một cuốn sách không đọc được là kiểu hỏng tệ
    nhất - không có lỗi nào báo ra, số chunk vẫn khác 0, và người dùng chỉ phát hiện khi câu
    trả lời sai. Với PDF scan (rất phổ biến với giáo trình tiếng Việt) thì đây không phải
    trường hợp hiếm mà là mặc định.

    Chỉ đếm chữ THẬT (loại bản ghi ảnh ra): một tài liệu có 400 chú thích ảnh mà không có
    dòng văn bản nào vẫn là một tài liệu không đọc được.
    """
    theo_nguon: Dict[str, int] = {}
    for m in cac_trang:
        if m.get("loai_noi_dung") != "anh":
            theo_nguon[m["nguon"]] = theo_nguon.get(m["nguon"], 0) + len(m["noidung"])

    for duong_dan in cac_duong_dan:
        so_ky_tu = theo_nguon.get(duong_dan.name, 0)
        if so_ky_tu >= config.SO_KY_TU_TOI_THIEU_MOT_TAI_LIEU:
            continue
        logger.error(
            "KHÔNG ĐỌC ĐƯỢC NỘI DUNG từ '%s' (chỉ %d ký tự văn bản). Gần như chắc chắn đây "
            "là PDF SCAN (ảnh chụp trang, không có lớp text). Hệ thống sẽ không trả lời "
            "đúng được về tài liệu này. Cách xử lý: bật OCR bằng BAT_OCR_DU_PHONG=1 và "
            "'ollama pull %s', rồi bấm Đọc tài liệu.",
            duong_dan.name, so_ky_tu, config.VISION_MODEL_NAME,
        )


def cac_file_tai_lieu(thu_muc: Path) -> List[Path]:
    """Danh sách file tài liệu hỗ trợ được trong thư mục, thứ tự ổn định.

    Tách ra thành hàm riêng vì luồng index tăng dần (app.py) cần đúng danh sách này để đối
    chiếu băm nội dung, mà không nên tự viết lại phép lọc phần đuôi file ở một chỗ thứ hai.
    """
    return [d for d in sorted(thu_muc.glob("*")) if d.suffix.lower() in CAC_DUOI_HO_TRO]


def doc_nhieu_file(cac_duong_dan: List[Path]) -> List[Dict]:
    """Đọc một DANH SÁCH file cụ thể (có cache), gom lỗi lại thay vì để một file hỏng
    làm sập cả lần build.

    Nhận danh sách thay vì thư mục là điều kiện để index tăng dần hoạt động: khi chỉ có 1
    trong 11 tài liệu thay đổi, app.py chỉ đưa đúng file đó vào đây.
    """
    ket_qua = []
    cac_file_loi = []
    for thu_tu, duong_dan in enumerate(cac_duong_dan, start=1):
        logger.info("[%d/%d] Đang xử lý '%s'...", thu_tu, len(cac_duong_dan), duong_dan.name)
        try:
            ket_qua.extend(doc_tai_lieu_co_cache(duong_dan))
        except Exception as loi:  # noqa: BLE001 - xem giải thích ngay dưới
            # MỘT FILE HỎNG KHÔNG ĐƯỢC LÀM SẬP CẢ LẦN BUILD. Trước đây một file PDF đặt mật
            # khẩu, tải về dở dang, hay dùng cấu trúc lạ sẽ ném lỗi và giết luôn toàn bộ luồng
            # Ingestion - người dùng mất trắng công build của 12 tài liệu còn lại và chỉ nhận
            # được một traceback. Với một hệ thống mà người dùng tự nạp tài liệu bất kỳ vào,
            # gặp file đọc không được là chuyện bình thường, không phải ngoại lệ.
            #
            # Bắt Exception rộng là có chủ đích: pdfplumber/python-docx/python-pptx ném ra rất
            # nhiều loại lỗi khác nhau tuỳ file hỏng kiểu gì, liệt kê hết là không thể. Cái
            # quan trọng là KHÔNG NUỐT LỖI - tên file và loại lỗi được ghi ra rõ ràng, và
            # được tổng kết lại ở cuối để người dùng biết chính xác tài liệu nào không vào
            # được index.
            cac_file_loi.append((duong_dan.name, f"{type(loi).__name__}: {loi}"))
            logger.error(
                "KHÔNG ĐỌC ĐƯỢC FILE '%s' (%s: %s) - bỏ qua file này và tiếp tục với các "
                "tài liệu còn lại. File có thể đặt mật khẩu, tải về dở dang, hoặc sai định "
                "dạng so với phần đuôi.",
                duong_dan.name, type(loi).__name__, loi,
            )

    if cac_file_loi:
        logger.error(
            "TỔNG KẾT: %d/%d tài liệu KHÔNG vào được index: %s",
            len(cac_file_loi),
            len(cac_file_loi) + len({m["nguon"] for m in ket_qua}),
            ", ".join(ten for ten, _ in cac_file_loi),
        )

    # Chú thích ảnh bằng model vision KHÔNG còn nằm ở đây mà đã chuyển vào
    # doc_tai_lieu_hoan_chinh() - tức làm xong trọn vẹn từng tài liệu một. Bản trước gom ảnh
    # của cả corpus lại rồi xử lý một lượt để báo được tiến độ "ảnh i/n"; nhưng khi kết quả
    # đọc được đem đi cache theo từng tài liệu thì một bản ghi ảnh CHƯA có chú thích là một
    # bản ghi chưa hoàn chỉnh - cache nó lại nghĩa là đóng băng một kết quả thiếu. Tiến độ
    # nay báo theo tài liệu, hợp với thực tế là đa số tài liệu không được xử lý lại.
    _canh_bao_tai_lieu_khong_doc_duoc(ket_qua, cac_duong_dan)
    return ket_qua


def doc_thu_muc(thu_muc: Path) -> List[Dict]:
    """Đọc toàn bộ file PDF/PPTX/DOCX trong 1 thư mục, dùng cho luồng Ingestion."""
    return doc_nhieu_file(cac_file_tai_lieu(thu_muc))
