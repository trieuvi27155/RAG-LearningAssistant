"""Test cho document_loader - đọc bảng, đệ quy group shape, nhận diện tiêu đề.

Tài liệu test được DỰNG TẠI CHỖ bằng chính python-docx/python-pptx trong thư mục tạm, thay
vì kèm file nhị phân vào repo: tránh phải bảo trì fixture nhị phân (không đọc/diff được),
và bản thân việc dựng file đã ghi rõ tài liệu chứa đúng những gì. Phần nhận diện tiêu đề
PDF được test trên dữ liệu ký tự tổng hợp - không cần file PDF thật, chạy trong mili giây.
"""

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import pytest
from docx import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

import config
from rag.document_loader import (
    MOC_BANG_MO,
    _bang_sang_markdown,
    _khoi_bang,
    _la_bang_that,
    _o_bang_khong_lap,
    _phat_hien_tieu_de_pdf,
    doc_docx,
    doc_pptx,
    duyet_shape,
)


class TrangGia:
    """Giả lập pdfplumber Page - chỉ cần .chars cho việc nhận diện tiêu đề."""

    def __init__(self, cac_dong):
        # cac_dong: list (text, co_chu, toa_do_top)
        self.chars = [
            {"text": ky_tu, "size": co_chu, "top": top}
            for text, co_chu, top in cac_dong
            for ky_tu in text
        ]


# ======================================================================
# Bảng -> Markdown
# ======================================================================

def test_bang_sang_markdown_giu_dung_hang_cot():
    md = _bang_sang_markdown([["Vai trò", "Mượn"], ["Thủ thư", "Có"], ["Khách", "Không"]])
    dong = md.splitlines()
    assert dong[0] == "| Vai trò | Mượn |"
    assert dong[1] == "| --- | --- |"
    assert "| Thủ thư | Có |" in dong


def test_bang_sang_markdown_bu_o_thieu_va_bo_hang_rong():
    """Bảng thật hay có hàng thiếu ô (ô gộp) hoặc hàng rỗng hoàn toàn - nếu không xử lý,
    số cột giữa các hàng lệch nhau và bảng Markdown vỡ."""
    md = _bang_sang_markdown([["A", "B", "C"], ["chỉ một ô"], [None, None, None], []])
    for dong in md.splitlines():
        assert dong.count("|") == 4  # 3 cột -> luôn 4 dấu gạch đứng
    assert "None" not in md


def test_bang_rong_tra_ve_chuoi_rong():
    assert _bang_sang_markdown([]) == ""
    assert _bang_sang_markdown([[None, None]]) == ""


# ======================================================================
# DOCX: bảng theo đúng thứ tự tài liệu + tiêu đề
# ======================================================================

def _tao_docx_tam(thu_muc: Path) -> Path:
    duong_dan = thu_muc / "thu.docx"
    tai_lieu = Document()
    tai_lieu.add_heading("Chương một", level=1)
    tai_lieu.add_paragraph("Đoạn văn giới thiệu trước bảng, đủ dài để không bị lọc bỏ.")
    bang = tai_lieu.add_table(rows=2, cols=2)
    bang.cell(0, 0).text = "Cột A"
    bang.cell(0, 1).text = "Cột B"
    bang.cell(1, 0).text = "Giá trị duy nhất XYZ"
    bang.cell(1, 1).text = "123"
    tai_lieu.add_paragraph("Đoạn văn kết thúc nằm sau bảng, cũng đủ dài để giữ lại.")
    tai_lieu.save(duong_dan)
    return duong_dan


def test_docx_doc_duoc_noi_dung_bang():
    """Bản trước chỉ duyệt document.paragraphs nên bảng biến mất hoàn toàn khỏi index."""
    with tempfile.TemporaryDirectory() as thu_muc:
        cac_trang = doc_docx(_tao_docx_tam(Path(thu_muc)))
    noi_dung = "\n".join(t["noidung"] for t in cac_trang)
    assert "Giá trị duy nhất XYZ" in noi_dung
    assert MOC_BANG_MO in noi_dung


def test_docx_bang_nam_dung_thu_tu_giua_hai_doan_van():
    """Bảng phải nằm CẠNH đoạn văn giới thiệu nó. Gom hết bảng về cuối trang sẽ khiến ngữ
    cảnh quanh bảng sai, dù nội dung bảng vẫn có mặt."""
    with tempfile.TemporaryDirectory() as thu_muc:
        noi_dung = doc_docx(_tao_docx_tam(Path(thu_muc)))[0]["noidung"]
    assert noi_dung.index("giới thiệu trước bảng") < noi_dung.index("Giá trị duy nhất XYZ")
    assert noi_dung.index("Giá trị duy nhất XYZ") < noi_dung.index("kết thúc nằm sau bảng")


def test_docx_danh_dau_tieu_de_theo_style_heading():
    with tempfile.TemporaryDirectory() as thu_muc:
        noi_dung = doc_docx(_tao_docx_tam(Path(thu_muc)))[0]["noidung"]
    assert "# Chương một" in noi_dung


# ======================================================================
# PPTX: đệ quy group shape
# ======================================================================

def _tao_pptx_co_group(thu_muc: Path) -> Path:
    duong_dan = thu_muc / "thu.pptx"
    trinh_chieu = Presentation()
    slide = trinh_chieu.slides.add_slide(trinh_chieu.slide_layouts[5])
    slide.shapes.title.text = "Tiêu đề slide thử nghiệm"
    nhom = slide.shapes.add_group_shape()
    khung = nhom.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
    khung.text_frame.text = "Nội dung ẩn trong nhóm hình ABC"
    trinh_chieu.save(duong_dan)
    return duong_dan


def test_pptx_doc_duoc_noi_dung_trong_group_shape():
    """Vòng lặp `for shape in slide.shapes` phẳng KHÔNG thấy gì bên trong group - nội dung
    biến mất khỏi index mà không có lỗi nào báo ra (đã đo trên tài liệu mẫu)."""
    with tempfile.TemporaryDirectory() as thu_muc:
        cac_trang = doc_pptx(_tao_pptx_co_group(Path(thu_muc)))
    assert "Nội dung ẩn trong nhóm hình ABC" in cac_trang[0]["noidung"]


def test_duyet_shape_de_quy_vao_nhom():
    """duyet_shape trả về các shape LÁ (đi xuyên qua nhóm), nên cái textbox nằm trong nhóm
    phải có mặt ở kết quả dù nó không hề xuất hiện ở tầng trên cùng."""
    with tempfile.TemporaryDirectory() as thu_muc:
        trinh_chieu = Presentation(_tao_pptx_co_group(Path(thu_muc)))
        slide = trinh_chieu.slides[0]

        def co_text_ben_trong(cac_shape):
            return any(
                s.has_text_frame and "nhóm hình ABC" in s.text_frame.text for s in cac_shape
            )

        assert not co_text_ben_trong(slide.shapes), "tầng trên cùng chỉ thấy cái nhóm"
        assert co_text_ben_trong(duyet_shape(slide.shapes)), "phải đi được vào trong nhóm"
        # Và không được trả về chính cái nhóm (nhóm không có nội dung của riêng nó).
        assert all(s.shape_type != MSO_SHAPE_TYPE.GROUP for s in duyet_shape(slide.shapes))


def test_pptx_tieu_de_slide_thanh_dau_tieu_de():
    with tempfile.TemporaryDirectory() as thu_muc:
        noi_dung = doc_pptx(_tao_pptx_co_group(Path(thu_muc)))[0]["noidung"]
    assert "## Tiêu đề slide thử nghiệm" in noi_dung
    # Không được xuất hiện 2 lần (1 lần là tiêu đề, 1 lần là text thường).
    assert noi_dung.count("Tiêu đề slide thử nghiệm") == 1


# ======================================================================
# PDF: nhận diện tiêu đề bằng cỡ chữ (heuristic)
# ======================================================================

def test_tieu_de_pdf_can_ca_chu_to_va_dong_ngan():
    """Hai điều kiện phải cùng đúng. Chỉ dựa vào cỡ chữ sẽ nhận nhầm cả đoạn văn in cỡ lớn;
    chỉ dựa vào độ dài sẽ nhận nhầm mọi dòng ngắn cuối đoạn."""
    # Cỡ chữ áp đảo được tính theo SỐ KÝ TỰ, nên phần thân phải thật sự chiếm đa số ký tự -
    # đúng như tài liệu thật (vài dòng tiêu đề giữa hàng nghìn ký tự thân bài).
    cac_dong = [("Chương một", 20.0, 10)]            # to + ngắn -> tiêu đề
    cac_dong += [
        (f"Đây là câu văn bình thường của phần thân, dòng số {i}.", 11.0, 40 + i * 20)
        for i in range(8)
    ]
    cac_dong.append(("Một dòng ngắn", 11.0, 400))    # ngắn nhưng KHÔNG to -> không phải
    cac_dong.append(("A" * 200, 20.0, 430))          # to nhưng QUÁ DÀI -> không phải
    tieu_de = _phat_hien_tieu_de_pdf(TrangGia(cac_dong))
    assert "Chương một" in tieu_de
    assert "Một dòng ngắn" not in tieu_de
    assert not any(len(t) > config.DO_DAI_TOI_DA_TIEU_DE for t in tieu_de)


def test_tieu_de_pdf_so_theo_ti_le_khong_theo_so_cung():
    """Mỗi tài liệu dùng cỡ chữ nền khác nhau (11pt, 13pt...). So theo TỈ LỆ với cỡ chữ áp
    đảo của trang mới tổng quát; so với một con số cứng sẽ sai ngay khi đổi tài liệu."""
    # Trang cỡ chữ nền LỚN (18pt): dòng 20pt chỉ hơn 1.11 lần -> chưa đủ ngưỡng 1.15.
    trang = TrangGia([("Không phải tiêu đề", 20.0, 10)] + [
        (f"Dòng thân bài số {i} trong tài liệu khổ lớn.", 18.0, 40 + i * 30) for i in range(5)
    ])
    assert "Không phải tiêu đề" not in _phat_hien_tieu_de_pdf(trang)


def test_trang_pdf_rong_khong_gay_loi():
    assert _phat_hien_tieu_de_pdf(TrangGia([])) == set()


# ======================================================================
# Bảng giả và ô gộp - các lỗi chỉ lộ ra trên tài liệu thật
# ======================================================================

def test_bang_mot_cot_khong_duoc_coi_la_bang():
    """pdfplumber dò bảng bằng đường kẻ, nên khung viền trang trí và gạch chân tiêu đề
    trong slide bị dò nhầm thành "bảng" 1 cột. Đo trên bộ slide Computer Vision thật: gần
    như mọi trang đều dính, và vì vùng bảng bị loại khỏi text thường nên TIÊU ĐỀ SLIDE -
    phần đậm đặc thông tin nhất - bị bóc khỏi văn bản."""
    assert not _la_bang_that([["COMPUTER VISION"], [""]])
    assert not _la_bang_that([["CONTENTS"], ["1. Detector and Descriptor"]])


def test_bang_that_phai_co_it_nhat_hai_hang_hai_cot():
    assert _la_bang_that([["Tên", "Giá"], ["Sách", "10.000"]])
    assert not _la_bang_that([["Tên", "Giá"]])  # chỉ 1 hàng -> không có quan hệ hàng-cột


def test_khoi_bang_bo_qua_bang_gia():
    assert _khoi_bang([["CHỈ MỘT TIÊU ĐỀ"], [""]]) == ""
    assert MOC_BANG_MO in _khoi_bang([["A", "B"], ["1", "2"]])


class _OGia:
    """Giả lập ô của python-docx: ô gộp trả về CÙNG một phần tử XML nền."""

    def __init__(self, text, nen):
        self.text = text
        self._tc = nen


class _HangGia:
    def __init__(self, cells):
        self.cells = cells


def test_khu_o_gop_bi_nhan_ban():
    """Biểu mẫu hành chính gộp ô rất nhiều; thư viện trả cùng một ô cho mọi cột mà ô gộp
    trải qua. Đo trên file đăng ký đề tài NCKH thật: 131.115 ký tự cho một biểu mẫu vài
    trang, tức nội dung bị nhân lên ~10 lần - vừa phình index vừa hỏng retrieval."""
    nen_chung = object()
    hang = _HangGia([
        _OGia("ĐĂNG KÝ ĐỀ TÀI", nen_chung),
        _OGia("ĐĂNG KÝ ĐỀ TÀI", nen_chung),
        _OGia("ĐĂNG KÝ ĐỀ TÀI", nen_chung),
        _OGia("Lĩnh vực", object()),
    ])
    assert _o_bang_khong_lap(hang) == ["ĐĂNG KÝ ĐỀ TÀI", "Lĩnh vực"]


def test_khong_xoa_nham_hai_o_khac_nhau_trung_noi_dung():
    """Khử trùng phải dựa vào phần tử XML nền, KHÔNG dựa vào text: hai ô độc lập hoàn toàn
    có thể tình cờ trùng nội dung ("Có"/"Có") - xoá đi là mất dữ liệu thật."""
    hang = _HangGia([_OGia("Có", object()), _OGia("Có", object())])
    assert _o_bang_khong_lap(hang) == ["Có", "Có"]


def test_o_rong_van_duoc_khu_trung_dung():
    """Phần tử lxml rỗng bị coi là falsy khi truth-test, nên viết `a or b` sẽ rơi nhầm
    nhánh và hỏng việc khử trùng - phải kiểm tra `is not None`."""
    class _NenRong:
        def __len__(self):
            return 0  # lxml: phần tử không có con -> falsy

    nen_chung = _NenRong()
    hang = _HangGia([_OGia("", nen_chung), _OGia("", nen_chung), _OGia("X", object())])
    assert _o_bang_khong_lap(hang) == ["", "X"]
